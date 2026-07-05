import os
import sys
import asyncio
import json
import re
import pydantic
import shutil
import copy
from typing import List, Literal, Optional, Dict, Any
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from google.genai import types
from playwright.sync_api import sync_playwright

if sys.platform == 'win32':
    try:
        asyncio.set_event_loop_policy(asyncio.WindowsProactorEventLoopPolicy())
    except Exception:
        pass

try:
    from gp_chat import state_manager
    from gp_chat import llm_router
except ImportError:
    import state_manager
    import llm_router

# --- Pydantic DSL スキーマ定義 ---


class CropAreaSchema(pydantic.BaseModel):
    ymin: int = pydantic.Field(..., description="切り出し領域の上端 (0-1000)")
    xmin: int = pydantic.Field(..., description="切り出し領域の左端 (0-1000)")
    ymax: int = pydantic.Field(..., description="切り出し領域の下端 (0-1000)")
    xmax: int = pydantic.Field(..., description="切り出し領域の右端 (0-1000)")


class PlaceholderContent(pydantic.BaseModel):
    idx: int = pydantic.Field(..., description="流し込み先のプレースホルダーインデックス (テンプレート定義のもの)。")
    text_content: Optional[str] = pydantic.Field(
        None,
        description="流し込むテキスト。箇条書きの場合は各要素の先頭に '• ' を付けず、改行 '\\n' で区切って複数行のテキストとして記述してください（PowerPoint側で自動的に箇条書きが適用されるため）。"
    )
    image_prompt: Optional[str] = pydantic.Field(
        None,
        description="このプレースホルダーが IMAGE タイプの場合に、AI画像生成 (Imagen 3) でイラストを自動生成させるための英語プロンプト。"
    )
    use_user_image: bool = pydantic.Field(
        False,
        description="ユーザー添付画像を使用する枠である場合は True。"
    )
    crop_instruction: Optional[str] = pydantic.Field(
        None,
        description="ユーザー添付画像を使用する際のトリミング指示。"
    )


class SlideNode(pydantic.BaseModel):
    slide_number: int
    title: str = pydantic.Field(..., description="スライドのタイトル。")
    layout_name: str = pydantic.Field(..., description="使用するテンプレートのスライドレイアウト名（動的に指示された一覧から正確に一致するものを選択）。")
    placeholders: List[PlaceholderContent] = pydantic.Field(..., description="プレースホルダーごとの流し込みデータの一覧。タイトル(TITLE/CENTER_TITLE)のプレースホルダー(idx=0等)にも必ずタイトル文字列を指定してください。")
    visual_type: Literal["auto", "none", "summary", "timeline", "process", "comparison", "kpi", "matrix", "risk"] = pydantic.Field(
        "auto",
        description="Optional generic visual treatment. Choose from the information structure, not the topic: timeline/process/comparison/kpi/matrix/risk/summary/none/auto.",
    )
    visual_variant: str = pydantic.Field(
        "auto",
        description=(
            "Optional visual variant. Examples: cards_2x2, scorecard, pros_cons, "
            "metric_cards, big_numbers, vertical_flow, chevron_flow, vertical_timeline, "
            "phase_bands, horizontal_timeline, milestone_cards, progress_bars, "
            "priority_quadrants, hub_spoke, risk_matrix, cause_impact_mitigation, "
            "bands, pyramid, auto."
        ),
    )
    coverage_refs: List[str] = pydantic.Field(
        default_factory=list,
        description=(
            "Exact source brief facts, source_coverage_units, or coverage_requirements "
            "this slide reflects. Used for content coverage auditing."
        ),
    )


class PresentationDSLSchema(pydantic.BaseModel):
    presentation_title: str = pydantic.Field(..., max_length=30)
    slides: List[SlideNode]


class PresentationCoverageAudit(pydantic.BaseModel):
    status: Literal["pass", "warning", "fail"]
    covered_items: List[str] = pydantic.Field(default_factory=list)
    weakly_covered_items: List[str] = pydantic.Field(default_factory=list)
    missing_items: List[str] = pydantic.Field(default_factory=list)
    overcompressed_items: List[str] = pydantic.Field(default_factory=list)
    recommendations: List[str] = pydantic.Field(default_factory=list)
    revision_instruction: str = ""


class SlideVisualAudit(pydantic.BaseModel):
    slide_number: int
    status: Literal["pass", "warning", "fail"]
    text_overflow: bool
    text_overlap: bool
    unreadable_text: bool
    excessive_blank_space: bool
    poor_image_relevance: bool
    template_artifacts_visible: bool
    issues: List[str] = pydantic.Field(default_factory=list)
    recommendations: List[str] = pydantic.Field(default_factory=list)
    overall_comment: str = ""


class ReferenceEntry(pydantic.BaseModel):
    reference_id: str = pydantic.Field("", description="出典ID。例: R1")
    source_type: str = pydantic.Field("", description="web / attachment / conversation / canvas / other")
    title: str = pydantic.Field("", description="出典名、ページタイトル、ファイル名、または資料名。")
    uri: Optional[str] = pydantic.Field(None, description="Web URL。添付や会話の場合は空でよい。")
    snippet: str = pydantic.Field("", description="根拠本文スニペット、または出典内容の短い要約。")
    used_slides: List[int] = pydantic.Field(default_factory=list, description="この出典を使った本文スライド番号。")


class PresentationSourceBrief(pydantic.BaseModel):
    core_request: str = pydantic.Field("", description="ユーザーが作りたい資料の目的・問い。")
    audience: str = pydantic.Field("", description="想定読者・利用場面。明示がなければ推定でよい。")
    source_inventory: List[str] = pydantic.Field(default_factory=list, description="会話、添付、画像、検索など利用可能な材料の一覧。")
    key_facts: List[str] = pydantic.Field(default_factory=list, description="スライドに反映すべき主要事実。日付・数値・固有名詞を優先。")
    evidence_notes: List[str] = pydantic.Field(default_factory=list, description="根拠、出典、添付資料、検索結果への参照メモ。")
    visual_assets: List[str] = pydantic.Field(default_factory=list, description="使う価値がある添付画像・図表・生成画像候補。不要なら空。")
    recommended_storyline: List[str] = pydantic.Field(default_factory=list, description="推奨スライド構成。")
    image_policy: str = pydantic.Field("", description="画像生成や添付画像利用の採否方針。")
    gaps_or_uncertainties: List[str] = pydantic.Field(default_factory=list, description="不明点、確度が低い点、資料に入れる際の注意。")
    coverage_requirements: List[str] = pydantic.Field(default_factory=list, description="最終PPTXに必ず反映すべき観点。")
    source_coverage_units: List[str] = pydantic.Field(
        default_factory=list,
        description=(
            "Atomic source facts or user-provided information units that should be traced "
            "into the final deck when relevant."
        ),
    )
    references: List[ReferenceEntry] = pydantic.Field(
        default_factory=list,
        description="PPTX末尾の出典スライドと詳細JSONに載せる参照情報。Web URL、添付名、根拠スニペットを含める。",
    )


def _normalize_api_role(role: str) -> str:
    if role in ("assistant", "model"):
        return "model"
    return "user"


def _clone_content_for_pptx(content):
    try:
        if hasattr(content, "model_copy"):
            return content.model_copy(deep=True)
        return copy.deepcopy(content)
    except Exception:
        return content


def _format_attachment_summary(file_attachments_meta: Optional[List[dict]]) -> str:
    if not file_attachments_meta:
        return ""
    lines = ["\n\n【添付ファイル情報】"]
    for item in file_attachments_meta:
        name = item.get("name", "unknown")
        file_type = item.get("type", "file")
        size = item.get("size")
        size_text = f", size={size}" if size is not None else ""
        lines.append(f"- {name} ({file_type}{size_text})")
    return "\n".join(lines)


def _format_grounding_metadata(metadata: Optional[dict]) -> str:
    if not metadata:
        return ""
    lines = []
    queries = metadata.get("queries") or []
    sources = metadata.get("sources") or []
    if queries:
        lines.append("検索クエリ:")
        for query in queries[:10]:
            lines.append(f"- {query}")
    if sources:
        lines.append("検索ソース:")
        for source in sources[:20]:
            title = source.get("title") or source.get("uri") or "source"
            uri = source.get("uri") or ""
            lines.append(f"- {title}: {uri}")
    return "\n".join(lines)


def _shorten_text(value: Any, max_chars: int = 180) -> str:
    text = re.sub(r"\s+", " ", str(value or "")).strip()
    if len(text) <= max_chars:
        return text
    return text[: max_chars - 1].rstrip() + "…"


def _reference_identity(reference: ReferenceEntry) -> str:
    uri = (reference.uri or "").strip().lower()
    if uri:
        return f"uri:{uri}"
    title = (reference.title or "").strip().lower()
    source_type = (reference.source_type or "").strip().lower()
    return f"{source_type}:{title}"


def _reference_from_source(source: dict, index: int) -> Optional[ReferenceEntry]:
    uri = source.get("uri") or source.get("url")
    title = source.get("title") or uri
    if not title and not uri:
        return None
    snippet = (
        source.get("snippet")
        or source.get("summary")
        or source.get("text")
        or "Web検索で参照された出典。"
    )
    return ReferenceEntry(
        reference_id=f"R{index}",
        source_type="web",
        title=_shorten_text(title, 120),
        uri=uri,
        snippet=_shorten_text(snippet, 220),
    )


def _reference_from_attachment(item: dict, index: int) -> Optional[ReferenceEntry]:
    name = item.get("name")
    if not name:
        return None
    file_type = item.get("type", "attachment")
    size = item.get("size")
    size_text = f", size={size}" if size is not None else ""
    return ReferenceEntry(
        reference_id=f"R{index}",
        source_type="attachment",
        title=_shorten_text(name, 120),
        uri=None,
        snippet=_shorten_text(f"添付ファイル ({file_type}{size_text}) としてLLM入力へ取り込まれた資料。", 220),
    )


def _finalize_reference_entries(
    references: List[ReferenceEntry],
    file_attachments_meta: Optional[List[dict]] = None,
    grounding_metadata: Optional[dict] = None,
) -> List[ReferenceEntry]:
    merged: List[ReferenceEntry] = []
    seen: set[str] = set()

    def add_reference(reference: Optional[ReferenceEntry]) -> None:
        if not reference:
            return
        title = _shorten_text(reference.title, 140)
        snippet = _shorten_text(reference.snippet, 240)
        if not title and not snippet:
            return
        normalized = ReferenceEntry(
            reference_id=reference.reference_id,
            source_type=(reference.source_type or "other").strip() or "other",
            title=title or "Untitled source",
            uri=(reference.uri or None),
            snippet=snippet or "参照情報として利用。",
            used_slides=sorted({int(num) for num in (reference.used_slides or []) if isinstance(num, int) or str(num).isdigit()}),
        )
        identity = _reference_identity(normalized)
        if identity in seen:
            return
        seen.add(identity)
        merged.append(normalized)
    for reference in references or []:
        add_reference(reference)
    next_index = len(merged) + 1
    for source in (grounding_metadata or {}).get("sources", []) or []:
        before = len(merged)
        add_reference(_reference_from_source(source, next_index))
        if len(merged) > before:
            next_index += 1
    for item in file_attachments_meta or []:
        before = len(merged)
        add_reference(_reference_from_attachment(item, next_index))
        if len(merged) > before:
            next_index += 1
    if not merged:
        merged.append(
            ReferenceEntry(
                reference_id="R1",
                source_type="conversation",
                title="会話履歴",
                snippet="ユーザー依頼とこれまでの会話内容を資料化の主要材料として使用。",
            )
        )
    for index, reference in enumerate(merged, start=1):
        reference.reference_id = f"R{index}"
    return merged[:30]


def _attach_reference_usage(brief: PresentationSourceBrief, presentation_data: PresentationDSLSchema) -> None:
    if not brief.references:
        return
    for reference in brief.references:
        used = set(reference.used_slides or [])
        ref_id = reference.reference_id
        title_terms = [term for term in re.split(r"[\s　:/|｜,，、。()（）\[\]【】]+", reference.title or "") if len(term) >= 4]
        for slide in presentation_data.slides:
            visible = _slide_visible_text(slide)
            trace_text = "\n".join(slide.coverage_refs or [])
            haystack = f"{visible}\n{trace_text}"
            if ref_id and ref_id in haystack:
                used.add(slide.slide_number)
                continue
            if title_terms and any(term in haystack for term in title_terms[:4]):
                used.add(slide.slide_number)
        reference.used_slides = sorted(used)


def _brief_to_text(brief: PresentationSourceBrief) -> str:
    return json.dumps(brief.model_dump(), ensure_ascii=False, indent=2)


def _slide_visible_text(slide: SlideNode) -> str:
    chunks = [slide.title or ""]
    for content in slide.placeholders:
        if content.text_content:
            chunks.append(content.text_content)
        if content.image_prompt:
            chunks.append(f"[image_prompt] {content.image_prompt}")
        if content.use_user_image:
            chunks.append("[user_image]")
    return "\n".join(chunk for chunk in chunks if chunk).strip()


def _slide_body_text(slide: SlideNode) -> str:
    chunks = []
    for content in slide.placeholders:
        if content.text_content:
            chunks.append(content.text_content)
        if content.image_prompt:
            chunks.append(f"[image_prompt] {content.image_prompt}")
        if content.use_user_image:
            chunks.append("[user_image]")
    return "\n".join(chunk for chunk in chunks if chunk).strip()


def _required_coverage_items(brief: PresentationSourceBrief) -> List[str]:
    return list(
        dict.fromkeys(
            item
            for item in (
                list(brief.coverage_requirements)
                + list(brief.source_coverage_units)
                + list(brief.key_facts[:20])
            )
            if item
        )
    )


def _minimum_body_slide_count(brief: PresentationSourceBrief) -> int:
    storyline_count = len(brief.recommended_storyline or [])
    storyline_body_count = max(0, storyline_count - 1) if storyline_count else 0
    unit_count = len(brief.source_coverage_units or [])
    unit_body_count = (unit_count + 1) // 2 if unit_count else 0
    coverage_body_count = len(brief.coverage_requirements or [])
    candidates = [unit_body_count, coverage_body_count]
    if storyline_body_count:
        candidates.append(storyline_body_count)
    return max(3, min(12, max(candidates or [0])))


def _low_content_slide_numbers(slides: List[SlideNode], min_body_chars: int = 80) -> List[int]:
    low_content = []
    for slide in slides:
        body_text = _slide_body_text(slide)
        has_visual = any(content.image_prompt or content.use_user_image for content in slide.placeholders)
        if len(body_text) < min_body_chars and not has_visual:
            low_content.append(slide.slide_number)
    return low_content


def _deck_visible_text(slides: List[SlideNode]) -> str:
    return "\n\n".join(
        f"Slide {slide.slide_number}: {slide.title}\n{_slide_body_text(slide)}"
        for slide in slides
    )


def _build_conversation_excerpt(chat_history: List[dict], max_chars: int = 6000) -> str:
    lines = []
    for message in chat_history:
        role = message.get("role", "unknown")
        if role == "system":
            continue
        content = str(message.get("content", "")).strip()
        if not content:
            continue
        lines.append(f"{role}: {content}")
    excerpt = "\n\n".join(lines)
    return excerpt[:max_chars]


def _audit_summary_to_text(audits: List[SlideVisualAudit]) -> str:
    lines = []
    for audit in audits:
        if audit.status == "pass":
            continue
        issue_text = "; ".join(audit.issues) or audit.overall_comment
        rec_text = "; ".join(audit.recommendations)
        lines.append(
            f"- slide {audit.slide_number}: status={audit.status}, "
            f"overflow={audit.text_overflow}, overlap={audit.text_overlap}, "
            f"unreadable={audit.unreadable_text}, blank={audit.excessive_blank_space}, "
            f"poor_image={audit.poor_image_relevance}, artifacts={audit.template_artifacts_visible}. "
            f"issues={issue_text}. recommendations={rec_text}"
        )
    return "\n".join(lines) or "No serious visual issues."

# --- テンプレートレイアウトの自動解析スキャン ---

TITLE_PLACEHOLDER_TYPES = {1, 3}
CONTENT_PLACEHOLDER_TYPES = {2, 7, 12, 18}
LLM_PLACEHOLDER_TYPES = TITLE_PLACEHOLDER_TYPES | CONTENT_PLACEHOLDER_TYPES | {4}

PLACEHOLDER_TYPE_NAMES = {
    1: "TITLE",
    2: "BODY",
    3: "CENTER_TITLE",
    4: "SUBTITLE",
    7: "CONTENT",
    12: "TABLE",
    18: "IMAGE",
}


def scan_template_layouts(template_path: str) -> Dict[str, Any]:
    """PowerPointテンプレートを解析し、LLMに渡すためのレイアウト辞書情報を構築する"""
    layouts = {}
    if not os.path.exists(template_path):
        return layouts
    try:
        prs = Presentation(template_path)
        for i, layout in enumerate(prs.slide_layouts):
            placeholders = []
            for ph in layout.placeholders:
                ph_type = ph.placeholder_format.type
                if ph_type not in LLM_PLACEHOLDER_TYPES:
                    continue
                ph_type_name = PLACEHOLDER_TYPE_NAMES.get(ph_type, "TEXT")
                placeholders.append({
                    "idx": ph.placeholder_format.idx,
                    "name": ph.name,
                    "type": ph_type_name,
                    "left_in": round(ph.left.inches, 3),
                    "top_in": round(ph.top.inches, 3),
                    "width_in": round(ph.width.inches, 3),
                    "height_in": round(ph.height.inches, 3)
                })
            if not any(ph["type"] in {"BODY", "CONTENT", "TABLE", "IMAGE"} for ph in placeholders):
                continue
            layouts[layout.name] = {
                "layout_index": i,
                "placeholders": placeholders
            }
    except Exception as e:
        state_manager.add_debug_log(f"[PPTXAgent] Error scanning template layouts: {e}", "warning")
    return layouts

# --- 画像処理ユーティリティの実装 ---


def generate_ai_image_with_nano_banana(client, prompt: str, output_path: str) -> bool:
    """gemini-3.1-flash-lite-image (Nano Banana 2 Lite) を用いてオンデマンドで画像を生成し保存する。404エラーの場合は imagen-3.0-generate-002 へフォールバック。"""
    try:
        state_manager.add_debug_log(f"[PPTXAgent] Generating AI image using Nano Banana 2 Lite. Prompt: {prompt}")
        try:
            response = client.models.generate_images(
                model="gemini-3.1-flash-lite-image",
                prompt=prompt,
                config=types.GenerateImagesConfig(
                    number_of_images=1,
                    output_mime_type="image/png",
                    aspect_ratio="16:9",
                )
            )
        except Exception as e:
            err_msg = str(e)
            if "404" in err_msg or "NOT_FOUND" in err_msg or "permission" in err_msg.lower():
                state_manager.add_debug_log(f"[PPTXAgent] Nano Banana 2 Lite not available. Falling back to imagen-3.0-generate-002. Error: {e}", "warning")
                response = client.models.generate_images(
                    model="imagen-3.0-generate-002",
                    prompt=prompt,
                    config=types.GenerateImagesConfig(
                        number_of_images=1,
                        output_mime_type="image/png",
                        aspect_ratio="16:9",
                    )
                )
            else:
                raise
        if response.generated_images:
            img_data = response.generated_images[0].image.image_bytes
            with open(output_path, "wb") as f:
                f.write(img_data)
            state_manager.add_debug_log(f"[PPTXAgent] AI image generated and saved to {output_path}")
            return True
        else:
            state_manager.add_debug_log("[PPTXAgent] Image generation returned no images.", "warning")
    except Exception as e:
        state_manager.add_debug_log(f"[PPTXAgent] Image generation pipeline failed: {e}", "warning")
    return False


def crop_user_image_with_llm(client, image_bytes: bytes, image_mime_type: str, instruction: str, output_path: str) -> bool:
    """Geminiマルチモーダルで画像内の対象物を検出し、Pillowでトリミングして保存する"""
    try:
        from PIL import Image
        import io
        state_manager.add_debug_log(f"[PPTXAgent] Cropping user image based on instruction: {instruction}")
        # 1. 座標の推論要求
        system_instruction = (
            "あなたは画像認識および情報トリミングのスペシャリストです。\n"
            "入力された画像の中から、ユーザーの指示内容（切り出したい対象）に最も一致する部分のバウンディングボックスを決定してください。\n"
            "画像全体のサイズを縦横 1000 と仮定した相対座標（0〜1000の整数）で、[ymin, xmin, ymax, xmax] の形で座標を出力してください。"
        )
        gen_config = types.GenerateContentConfig(
            system_instruction=system_instruction,
            temperature=0.1,
            response_mime_type="application/json",
            response_schema=CropAreaSchema,
        )
        image_part = types.Part.from_bytes(data=image_bytes, mime_type=image_mime_type)
        user_prompt = f"指示: 「{instruction}」\n上記指示に合致する要素の座標範囲 [ymin, xmin, ymax, xmax] をJSONで出力してください。"
        result = client.models.generate_content(
            model="gemini-3.5-flash",
            contents=[image_part, user_prompt],
            config=gen_config
        )
        coords = CropAreaSchema.model_validate_json(result.text)
        state_manager.add_debug_log(f"[PPTXAgent] LLM detected crop area: {coords}")
        # 2. Pillowによる物理トリミング
        img = Image.open(io.BytesIO(image_bytes))
        w, h = img.size
        # 境界値のクリッピング
        left = int(max(0, min(1000, coords.xmin)) / 1000 * w)
        top = int(max(0, min(1000, coords.ymin)) / 1000 * h)
        right = int(max(0, min(1000, coords.xmax)) / 1000 * w)
        bottom = int(max(0, min(1000, coords.ymax)) / 1000 * h)
        # 最低サイズ保証
        if right - left < 10: right = min(w, left + 10)
        if bottom - top < 10: bottom = min(h, top + 10)
        cropped_img = img.crop((left, top, right, bottom))
        cropped_img.save(output_path)
        state_manager.add_debug_log(f"[PPTXAgent] Image cropped and saved: {output_path}")
        return True
    except Exception as e:
        state_manager.add_debug_log(f"[PPTXAgent] User image cropping failed: {e}", "warning")
        return False

# --- 生成済みPPTXの視覚品質検査 ---


def export_pptx_to_slide_images(pptx_path: str, output_dir: str) -> List[str]:
    """PowerPointで生成済みPPTXをPNG化し、スライド画像パスを返す。"""
    if sys.platform != "win32":
        state_manager.add_debug_log("[PPTXAgent] Visual audit skipped: PPTX image export requires Windows.", "warning")
        return []
    try:
        import glob
        import re
        import pythoncom
        import win32com.client
    except Exception as e:
        state_manager.add_debug_log(f"[PPTXAgent] Visual audit skipped: pywin32 unavailable: {e}", "warning")
        return []
    os.makedirs(output_dir, exist_ok=True)
    ppt_app = None
    presentation = None
    try:
        pythoncom.CoInitialize()
        ppt_app = win32com.client.Dispatch("PowerPoint.Application")
        presentation = ppt_app.Presentations.Open(os.path.abspath(pptx_path), ReadOnly=True, WithWindow=False)
        presentation.SaveAs(os.path.abspath(os.path.join(output_dir, "slide.png")), 18)
    except Exception as e:
        state_manager.add_debug_log(f"[PPTXAgent] Visual audit export failed: {e}", "warning")
        return []
    finally:
        if presentation:
            try:
                presentation.Close()
            except Exception:
                pass
        if ppt_app:
            try:
                ppt_app.Quit()
            except Exception:
                pass
        try:
            pythoncom.CoUninitialize()
        except Exception:
            pass
    image_paths = glob.glob(os.path.join(output_dir, "*.PNG")) or glob.glob(os.path.join(output_dir, "*.png"))
    nested_dir = os.path.join(output_dir, "slide")
    if not image_paths and os.path.isdir(nested_dir):
        image_paths = glob.glob(os.path.join(nested_dir, "*.PNG")) or glob.glob(os.path.join(nested_dir, "*.png"))

    def slide_sort_key(path: str):
        name = os.path.basename(path)
        nums = re.findall(r"\d+", name)
        return int(nums[-1]) if nums else 0
    return sorted(image_paths, key=slide_sort_key)


def run_visual_quality_audit(client, pptx_path: str, temp_dir: str, model_id: str = "gemini-3.5-flash") -> List[SlideVisualAudit]:
    """生成済みPPTXを画像化し、VLMで視覚品質を検査する。"""
    audit_dir = os.path.join(temp_dir, "visual_audit")
    image_paths = export_pptx_to_slide_images(pptx_path, audit_dir)
    if not image_paths:
        return []
    audit_model = model_id if "gemini" in model_id else "gemini-3.5-flash"
    audits: List[SlideVisualAudit] = []
    prompt = (
        "You are a strict PowerPoint visual QA reviewer. Inspect this slide image only.\n"
        "Evaluate whether the slide is professionally readable and properly laid out.\n"
        "Check for: clipped text, text overflowing outside placeholders, overlaps, unreadable small text, "
        "excessive blank space, irrelevant/generated-looking images, visible template placeholders/default text, "
        "and poor information hierarchy.\n"
        "Return JSON that exactly follows the schema. Use status='fail' for serious visual defects, "
        "status='warning' for noticeable but acceptable issues, and status='pass' only when the slide looks production-ready."
    )
    config = types.GenerateContentConfig(
        temperature=0.0,
        response_mime_type="application/json",
        response_schema=SlideVisualAudit,
    )
    for index, image_path in enumerate(image_paths, start=1):
        try:
            with open(image_path, "rb") as f:
                image_bytes = f.read()
            result = client.models.generate_content(
                model=audit_model,
                contents=[
                    types.Part.from_bytes(data=image_bytes, mime_type="image/png"),
                    types.Part.from_text(text=f"{prompt}\nSlide number: {index}"),
                ],
                config=config,
            )
            audit = SlideVisualAudit.model_validate_json(result.text)
            audit.slide_number = index
            audits.append(audit)
            if audit.status != "pass":
                state_manager.add_debug_log(
                    f"[PPTXAgent][VisualAudit] Slide {index}: {audit.status}. "
                    f"Issues: {'; '.join(audit.issues[:3]) or audit.overall_comment}",
                    "warning" if audit.status == "warning" else "error",
                )
        except Exception as e:
            state_manager.add_debug_log(f"[PPTXAgent] Visual audit failed on slide {index}: {e}", "warning")
    if audits:
        failed = sum(1 for item in audits if item.status == "fail")
        warned = sum(1 for item in audits if item.status == "warning")
        state_manager.add_debug_log(
            f"[PPTXAgent][VisualAudit] Completed: {len(audits)} slides checked, {failed} fail, {warned} warning."
        )
    return audits

# --- HTML/CSS テンプレート生成ヘルパー (幾何学検証用) ---


def generate_slide_html(slide: SlideNode, layouts_info: Dict[str, Any], font_size_offset: int = 0) -> str:
    """テンプレートのレイアウト実座標から幾何学検証用HTMLを動的に生成する"""
    def font_size_for_placeholder(ph_type: str) -> float:
        if ph_type in ("TITLE", "CENTER_TITLE"):
            return 18
        if ph_type in ("BODY", "SUBTITLE"):
            return 15
        if ph_type in ("CONTENT", "TABLE"):
            return 13.5
        return 12
    layout = layouts_info.get(slide.layout_name)
    layout_html = ""
    if layout:
        for ph in layout["placeholders"]:
            idx = ph["idx"]
            left_in = ph["left_in"]
            top_in = ph["top_in"]
            width_in = ph["width_in"]
            height_in = ph["height_in"]
            ph_type = ph.get("type", "CONTENT")
            ph_fs = max(9, font_size_for_placeholder(ph_type) + font_size_offset)
            # 対応するコンテンツをLLM出力から探す
            content_text = ""
            for content in slide.placeholders:
                if content.idx == idx:
                    if content.text_content:
                        bullets = content.text_content.split("\n")
                        content_text = "<ul>" + "".join(f"<li>{b}</li>" for b in bullets if b) + "</ul>"
                    break
            layout_html += f"""
            <div class="dynamic-text-block" id="ph_{idx}_{slide.slide_number}" style="position: absolute; left: {left_in}in; top: {top_in}in; width: {width_in}in; height: {height_in}in; font-size: {ph_fs}pt; overflow: hidden; line-height: 1.35; color: #374151;">
              {content_text}
            </div>
            """
    else:
        # テンプレートがない場合のシンプルなフォールバック
        bullet_fs = max(9, 18 + font_size_offset)
        layout_html = f"""
        <div class="dynamic-text-block" id="fallback_bullets_{slide.slide_number}" style="position: absolute; left: 1.0in; top: 2.2in; width: 11.333in; height: 4.0in; font-size: {bullet_fs}pt; overflow: hidden; line-height: 1.5; color: #374151;">
          {slide.title}
        </div>
        """
    html = f"""<!DOCTYPE html>
<html>
<head>
  <meta charset="utf-8">
  <style>
    body {{
      margin: 0;
      padding: 0;
      background: #ffffff;
      font-family: 'Meiryo', 'Malgun Gothic', sans-serif;
    }}
    .slide-container {{
      width: 13.333in;
      height: 7.5in;
      position: relative;
      background: #ffffff;
      box-sizing: border-box;
      overflow: hidden;
    }}
    .title {{
      position: absolute;
      left: 1.0in;
      top: 0.6in;
      width: 11.333in;
      height: 1.0in;
      font-size: 18pt;
      font-weight: bold;
      color: #111827;
      margin: 0;
      padding: 0;
      line-height: 1.2;
      display: flex;
      align-items: center;
    }}
    .dynamic-text-block {{
      box-sizing: border-box;
      overflow: hidden;
      word-wrap: break-word;
    }}
    ul {{
      margin: 0;
      padding-left: 20pt;
    }}
    li {{
      margin-bottom: 10pt;
    }}
  </style>
</head>
<body>
  <div class="slide-container">
    <div class="title dynamic-text-block" id="title_{slide.slide_number}">{slide.title}</div>
    {layout_html}
  </div>
</body>
</html>
"""
    return html

# --- PPTX 物理レンダリングヘルパー ---


def _format_text_placeholder(shape, font_size_offset: int = 0):
    if not shape.has_text_frame:
        return
    ph_type = shape.placeholder_format.type if shape.is_placeholder else None
    idx = shape.placeholder_format.idx if shape.is_placeholder else None
    if ph_type in (1, 3):  # TITLE / CENTER_TITLE
        base_size = 18
        bold = True
    elif ph_type == 2 or idx == 13:  # BODY / summary
        base_size = 15
        bold = True
    elif ph_type in (7, 12):  # CONTENT / TABLE
        base_size = 13.5
        bold = False
    elif ph_type == 4:  # SUBTITLE
        base_size = 15
        bold = False
    else:
        base_size = 12
        bold = False
    if (shape.text or "").lstrip().startswith(("※", "注", "Note")):
        base_size = 12
        bold = False
    font_size = max(8, base_size + font_size_offset)
    for paragraph in shape.text_frame.paragraphs:
        paragraph.font.name = "Meiryo"
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        if ph_type in (1, 3):
            paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.name = "Meiryo"
            run.font.size = Pt(font_size)
            run.font.bold = bold


def _format_cover_title(shape):
    if not shape.has_text_frame:
        return
    text_len = len(shape.text or "")
    if text_len <= 18:
        base_size = 34
    elif text_len <= 28:
        base_size = 30
    else:
        base_size = 26
    shape.text_frame.word_wrap = True
    for paragraph in shape.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.name = "Meiryo"
        paragraph.font.size = Pt(base_size)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(0x1e, 0x3a, 0x8a)
        for run in paragraph.runs:
            run.font.name = "Meiryo"
            run.font.size = Pt(base_size)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1e, 0x3a, 0x8a)


def _title_box_from_layout(slide, prs):
    for placeholder in slide.slide_layout.placeholders:
        if placeholder.placeholder_format.type in (1, 3):
            left = placeholder.left
            top = placeholder.top
            width = min(placeholder.width, prs.slide_width - left - Inches(0.3))
            height = min(placeholder.height, prs.slide_height - top - Inches(0.3))
            return left, top, width, height
    return Inches(0.6), Inches(2.0), prs.slide_width - Inches(1.2), Inches(1.2)

INFO_GREEN = RGBColor(0x2f, 0x6b, 0x1f)
INFO_GREEN_LIGHT = RGBColor(0xea, 0xf4, 0xe2)
INFO_BLUE = RGBColor(0x1e, 0x3a, 0x8a)
INFO_BLUE_LIGHT = RGBColor(0xe8, 0xee, 0xfb)
INFO_AMBER = RGBColor(0xf5, 0x9e, 0x0b)
INFO_AMBER_LIGHT = RGBColor(0xff, 0xf7, 0xe6)
INFO_GRAY = RGBColor(0x4b, 0x55, 0x63)
INFO_RED = RGBColor(0xdc, 0x26, 0x26)
INFO_RED_LIGHT = RGBColor(0xfe, 0xe2, 0xe2)
INFO_PURPLE = RGBColor(0x7c, 0x3a, 0xed)
INFO_PURPLE_LIGHT = RGBColor(0xf3, 0xe8, 0xff)

INFO_COLORS = [
    (INFO_BLUE_LIGHT, INFO_BLUE),
    (INFO_GREEN_LIGHT, INFO_GREEN),
    (INFO_AMBER_LIGHT, INFO_AMBER),
    (INFO_PURPLE_LIGHT, INFO_PURPLE),
    (INFO_RED_LIGHT, INFO_RED),
    (RGBColor(0xf8, 0xfa, 0xfc), RGBColor(0xcb, 0xd5, 0xe1)),
]

VISUAL_VARIANTS_BY_STYLE = {
    "timeline": (
        "vertical_timeline",
        "phase_bands",
        "horizontal_timeline",
        "milestone_cards",
        "now_next_later",
        "gantt_roadmap",
    ),
    "process": (
        "vertical_flow",
        "chevron_flow",
        "numbered_steps",
        "loop_cycle",
        "swimlane_flow",
        "funnel",
    ),
    "comparison": (
        "cards_2x2",
        "scorecard",
        "pros_cons",
        "table_compare",
        "ranked_bars",
        "before_after",
        "option_columns",
    ),
    "kpi": (
        "metric_cards",
        "big_numbers",
        "scorecard",
        "progress_bars",
        "gauge_cards",
        "delta_callouts",
        "waterfall",
    ),
    "risk": (
        "risk_matrix",
        "cause_impact_mitigation",
        "risk_register",
        "heatmap",
        "escalation_ladder",
    ),
    "matrix": (
        "risk_matrix",
        "priority_quadrants",
        "decision_matrix",
    ),
    "summary": (
        "bands",
        "pyramid",
        "hub_spoke",
    ),
}
VALID_VISUAL_VARIANTS = {
    variant for variants in VISUAL_VARIANTS_BY_STYLE.values() for variant in variants
}


def _set_shape_text(shape, text: str, font_size: float = 11, bold: bool = False, color: RGBColor = INFO_GRAY, align=PP_ALIGN.CENTER):
    shape.text_frame.clear()
    shape.text_frame.word_wrap = True
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = "Meiryo"
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    for run in p.runs:
        run.font.name = "Meiryo"
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color


def _add_info_box(slide, left, top, width, height, text, fill_color, line_color=INFO_GREEN, font_size=11, bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = Pt(1)
    _set_shape_text(shape, text, font_size=font_size, bold=bold)
    return shape


def _remove_placeholders_by_idx(slide, indexes: set[int]):
    for shape in list(slide.shapes):
        if shape.is_placeholder and shape.placeholder_format.idx in indexes:
            sp = shape._element
            sp.getparent().remove(sp)


def _placeholder_text(slide_data: SlideNode, indexes: set[int]) -> str:
    chunks = []
    for content in slide_data.placeholders:
        if content.idx in indexes and content.text_content:
            chunks.append(content.text_content)
    return "\n".join(chunks)


def _split_visual_items(text: str, limit: int = 6) -> List[str]:
    normalized = (
        text.replace(" | ", "\n")
        .replace("；", "\n")
        .replace(";", "\n")
        .replace("。", "\n")
        .replace("•", "\n")
        .replace("●", "\n")
        .replace("■", "\n")
    )
    raw_items = [item.strip(" -*－–—・•●■\t\r\n") for item in normalized.split("\n")]
    return [item for item in raw_items if item][:limit]


def _infer_visual_style(slide_data: SlideNode) -> str:
    requested = getattr(slide_data, "visual_type", "auto") or "auto"
    if requested != "auto":
        return requested
    text = (slide_data.title + "\n" + "\n".join(
        content.text_content or "" for content in slide_data.placeholders
    )).lower()
    items = _split_visual_items(text, limit=12)
    date_hits = re.findall(
        r"(?:\d{4}年\d{1,2}月|\d{4}[/-]\d{1,2}|\d{4}年|\d{4}[/-]|\d{1,2}[/-]\d{1,2}|q[1-4]|第[1-4]四半期|\d{1,2}月)",
        text,
        flags=re.IGNORECASE,
    )
    number_hits = re.findall(
        r"\d+(?:[.,]\d+)*(?:%|％|円|万円|億円|ドル|\$|件|人|社|日|時間|分|秒|ms|gb|mb|tb|トークン|token)?",
        text,
        flags=re.IGNORECASE,
    )
    colon_lines = [item for item in items if ":" in item or "：" in item]
    timeline_words = ("timeline", "milestone", "phase", "roadmap", "時系列", "工程", "フェーズ", "段階", "推移")
    process_words = ("step", "process", "flow", "input", "output", "workflow", "手順", "処理", "入力", "出力", "承認", "実行")
    risk_words = ("risk", "issue", "impact", "cause", "mitigation", "countermeasure", "リスク", "課題", "影響", "原因", "対策", "懸念")
    matrix_words = ("matrix", "quadrant", "priority", "importance", "urgency", "マトリクス", "重要度", "緊急度", "優先度")
    compare_words = ("compare", "comparison", "versus", " vs ", "差分", "比較", "対比", "選定", "評価軸")
    if len(date_hits) >= 2 or any(word in text for word in timeline_words):
        return "timeline"
    if any(word in text for word in matrix_words):
        return "matrix"
    if any(word in text for word in risk_words):
        return "risk"
    if len(colon_lines) >= 3 or any(word in text for word in compare_words):
        return "comparison"
    if len(number_hits) >= 3:
        return "kpi"
    if any(word in text for word in process_words):
        return "process"
    return "summary"


def _delete_shape(shape) -> None:
    sp = shape._element
    sp.getparent().remove(sp)


def _visual_placeholder_candidate(slide, slide_data: SlideNode):
    used_text_idxs = {
        content.idx
        for content in slide_data.placeholders
        if content.text_content or content.image_prompt or content.use_user_image
    }
    candidates = []
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
        ph_type = shape.placeholder_format.type
        idx = shape.placeholder_format.idx
        if idx not in used_text_idxs:
            continue
        if ph_type not in (7, 12, 18):
            continue
        area = shape.width * shape.height
        score = area + shape.left * 3 + shape.top
        if ph_type == 18:
            score += area * 2
        candidates.append((score, shape))
    if not candidates:
        return None
    candidates.sort(key=lambda item: item[0], reverse=True)
    return candidates[0][1]


def _draw_timeline(slide, items: List[str], left, top, width, height):
    line_x = left + Inches(0.18)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, line_x, top + Inches(0.12), Inches(0.03), height - Inches(0.24))
    line.fill.solid()
    line.fill.fore_color.rgb = INFO_GREEN
    line.line.fill.background()
    row_count = min(max(len(items), 2), 4)
    row_h = height / row_count
    for idx in range(row_count):
        item = items[idx] if idx < len(items) else ""
        y = top + row_h * idx + Inches(0.04)
        marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, line_x - Inches(0.07), y + Inches(0.16), Inches(0.17), Inches(0.17))
        marker.fill.solid()
        marker.fill.fore_color.rgb = INFO_GREEN
        marker.line.fill.background()
        _add_info_box(slide, left + Inches(0.38), y, width - Inches(0.44), row_h - Inches(0.08), item, INFO_GREEN_LIGHT, font_size=8.5, bold=idx == 0)


def _draw_kpi_grid(slide, items: List[str], style: str, left, top, width, height):
    cols = 2
    rows = 2
    gap = Inches(0.12)
    box_w = (width - gap) / cols
    box_h = (height - gap) / rows
    colors = [
        (INFO_BLUE_LIGHT, INFO_BLUE),
        (INFO_GREEN_LIGHT, INFO_GREEN),
        (INFO_AMBER_LIGHT, INFO_AMBER),
        (RGBColor(0xf8, 0xfa, 0xfc), RGBColor(0xcb, 0xd5, 0xe1)),
    ]
    for idx, item in enumerate(items[: cols * rows]):
        fill, line = colors[idx % len(colors)]
        col = idx % cols
        row = idx // cols
        _add_info_box(slide, left + (box_w + gap) * col, top + (box_h + gap) * row, box_w, box_h, item, fill, line_color=line, font_size=8.6, bold=True)


def _draw_big_numbers(slide, items: List[str], left, top, width, height):
    cols = min(max(len(items), 2), 3)
    gap = Inches(0.12)
    box_w = (width - gap * (cols - 1)) / cols
    for idx, item in enumerate(items[:cols]):
        value_match = re.search(r"(\d+(?:[.,]\d+)*(?:%|％|円|万円|億円|ドル|\$|件|人|社|日|時間|分|秒|ms|gb|mb|tb)?)", item, flags=re.IGNORECASE)
        value = value_match.group(1) if value_match else item.split(":")[0].split("：")[0]
        label = item.replace(value, "", 1).strip(" :：-")
        x = left + (box_w + gap) * idx
        box = _add_info_box(slide, x, top, box_w, height, "", RGBColor(0xf8, 0xfa, 0xfc), line_color=RGBColor(0xcb, 0xd5, 0xe1))
        box.text_frame.clear()
        p_value = box.text_frame.paragraphs[0]
        p_value.text = value[:16]
        p_value.alignment = PP_ALIGN.CENTER
        p_value.font.name = "Meiryo"
        p_value.font.size = Pt(16)
        p_value.font.bold = True
        p_value.font.color.rgb = INFO_BLUE
        p_label = box.text_frame.add_paragraph()
        p_label.text = label[:34]
        p_label.alignment = PP_ALIGN.CENTER
        p_label.font.name = "Meiryo"
        p_label.font.size = Pt(7.8)
        p_label.font.color.rgb = INFO_GRAY


def _draw_scorecard(slide, items: List[str], left, top, width, height):
    row_count = min(len(items), 5)
    if row_count <= 0:
        return
    row_h = height / row_count
    for idx, item in enumerate(items[:row_count]):
        y = top + row_h * idx
        fill = INFO_GREEN_LIGHT if idx % 2 == 0 else RGBColor(0xf8, 0xfa, 0xfc)
        _add_info_box(slide, left, y + Inches(0.03), width, row_h - Inches(0.06), item, fill, line_color=RGBColor(0xcb, 0xd5, 0xe1), font_size=8.4, bold=idx == 0)


def _set_list_text(shape, items: List[str], font_size: float = 8.4):
    shape.text_frame.clear()
    shape.text_frame.word_wrap = True
    shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    shape.text_frame.margin_top = Inches(0.12)
    shape.text_frame.margin_left = Inches(0.12)
    shape.text_frame.margin_right = Inches(0.08)
    for idx, item in enumerate(items):
        paragraph = shape.text_frame.paragraphs[0] if idx == 0 else shape.text_frame.add_paragraph()
        paragraph.text = f"- {item}"
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.font.name = "Meiryo"
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = INFO_GRAY
        for run in paragraph.runs:
            run.font.name = "Meiryo"
            run.font.size = Pt(font_size)
            run.font.color.rgb = INFO_GRAY


def _draw_pros_cons(slide, items: List[str], left, top, width, height):
    mid = left + width / 2
    gap = Inches(0.12)
    col_w = width / 2 - gap
    left_items = items[0::2][:3]
    right_items = items[1::2][:3] or items[3:6]
    _add_info_box(slide, left, top, col_w, Inches(0.34), "Pros", INFO_GREEN_LIGHT, line_color=INFO_GREEN, font_size=8.6, bold=True)
    _add_info_box(slide, mid + gap, top, col_w, Inches(0.34), "Cons", INFO_AMBER_LIGHT, line_color=INFO_AMBER, font_size=8.6, bold=True)
    left_box = _add_info_box(slide, left, top + Inches(0.44), col_w, height - Inches(0.44), "", RGBColor(0xf8, 0xfa, 0xfc), line_color=INFO_GREEN, font_size=8.4)
    right_box = _add_info_box(slide, mid + gap, top + Inches(0.44), col_w, height - Inches(0.44), "", RGBColor(0xf8, 0xfa, 0xfc), line_color=INFO_AMBER, font_size=8.4)
    _set_list_text(left_box, left_items, font_size=9.0)
    _set_list_text(right_box, right_items, font_size=9.0)


def _draw_flow(slide, items: List[str], left, top, width, height):
    step_count = min(max(len(items), 2), 4)
    box_h = (height - Inches(0.22) * (step_count - 1)) / step_count
    for idx in range(step_count):
        item = items[idx] if idx < len(items) else ""
        y = top + (box_h + Inches(0.22)) * idx
        _add_info_box(slide, left, y, width, box_h, item, INFO_BLUE_LIGHT if idx != 1 else INFO_AMBER_LIGHT, line_color=INFO_BLUE if idx != 1 else INFO_AMBER, font_size=8.5, bold=True)
        if idx < step_count - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left + width / 2 - Inches(0.08), y + box_h - Inches(0.02), Inches(0.16), Inches(0.22))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = INFO_GREEN
            arrow.line.fill.background()


def _draw_chevron_flow(slide, items: List[str], left, top, width, height):
    count = min(max(len(items), 2), 4)
    gap = Inches(0.05)
    box_w = (width - gap * (count - 1)) / count
    box_h = min(height, Inches(0.9))
    y = top + (height - box_h) / 2
    for idx in range(count):
        item = items[idx] if idx < len(items) else ""
        shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + (box_w + gap) * idx, y, box_w, box_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = [INFO_BLUE_LIGHT, INFO_GREEN_LIGHT, INFO_AMBER_LIGHT, RGBColor(0xf8, 0xfa, 0xfc)][idx % 4]
        shape.line.color.rgb = [INFO_BLUE, INFO_GREEN, INFO_AMBER, RGBColor(0xcb, 0xd5, 0xe1)][idx % 4]
        _set_shape_text(shape, item, font_size=7.8, bold=True)


def _draw_phase_bands(slide, items: List[str], left, top, width, height):
    count = min(max(len(items), 2), 4)
    band_h = height / count
    colors = [INFO_BLUE_LIGHT, INFO_GREEN_LIGHT, INFO_AMBER_LIGHT, RGBColor(0xf8, 0xfa, 0xfc)]
    lines = [INFO_BLUE, INFO_GREEN, INFO_AMBER, RGBColor(0xcb, 0xd5, 0xe1)]
    for idx in range(count):
        item = items[idx] if idx < len(items) else ""
        _add_info_box(slide, left + Inches(0.12) * idx, top + band_h * idx, width - Inches(0.12) * idx, band_h - Inches(0.08), item, colors[idx % 4], line_color=lines[idx % 4], font_size=8.4, bold=idx == 0)


def _draw_matrix(slide, items: List[str], left, top, width, height):
    mid_x = left + width / 2
    mid_y = top + height / 2
    h_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, mid_y, width, Inches(0.02))
    v_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, mid_x, top, Inches(0.02), height)
    for line in (h_line, v_line):
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0xcb, 0xd5, 0xe1)
        line.line.fill.background()
    box_w = width / 2 - Inches(0.12)
    box_h = height / 2 - Inches(0.12)
    positions = [
        (left, top, INFO_BLUE_LIGHT, INFO_BLUE),
        (mid_x + Inches(0.12), top, INFO_GREEN_LIGHT, INFO_GREEN),
        (left, mid_y + Inches(0.12), INFO_AMBER_LIGHT, INFO_AMBER),
        (mid_x + Inches(0.12), mid_y + Inches(0.12), RGBColor(0xf8, 0xfa, 0xfc), RGBColor(0xcb, 0xd5, 0xe1)),
    ]
    for idx, item in enumerate(items[:4]):
        x, y, fill, line = positions[idx]
        _add_info_box(slide, x, y, box_w, box_h, item, fill, line_color=line, font_size=8.0, bold=idx == 0)


def _draw_cause_impact_mitigation(slide, items: List[str], left, top, width, height):
    labels = ["Cause", "Impact", "Mitigation"]
    count = 3
    box_h = (height - Inches(0.24) * (count - 1)) / count
    for idx in range(count):
        text = items[idx] if idx < len(items) else labels[idx]
        y = top + (box_h + Inches(0.24)) * idx
        _add_info_box(slide, left, y, Inches(0.82), box_h, labels[idx], [INFO_BLUE_LIGHT, INFO_AMBER_LIGHT, INFO_GREEN_LIGHT][idx], line_color=[INFO_BLUE, INFO_AMBER, INFO_GREEN][idx], font_size=7.2, bold=True)
        _add_info_box(slide, left + Inches(0.92), y, width - Inches(0.92), box_h, text, RGBColor(0xf8, 0xfa, 0xfc), line_color=RGBColor(0xcb, 0xd5, 0xe1), font_size=8.0)


def _draw_summary_bands(slide, items: List[str], left, top, width, height):
    row_count = min(len(items), 4)
    if row_count <= 0:
        return
    row_h = min(Inches(0.58), height / row_count - Inches(0.04))
    for idx, item in enumerate(items[:4]):
        _add_info_box(
            slide,
            left,
            top + (row_h + Inches(0.1)) * idx,
            width,
            row_h,
            item,
            INFO_GREEN_LIGHT if idx % 2 == 0 else RGBColor(0xf8, 0xfa, 0xfc),
            line_color=INFO_GREEN if idx % 2 == 0 else RGBColor(0xcb, 0xd5, 0xe1),
            font_size=8.6,
            bold=idx == 0,
        )


def _draw_pyramid(slide, items: List[str], left, top, width, height):
    count = min(max(len(items), 2), 4)
    level_h = height / count
    for idx in range(count):
        shrink = (count - idx - 1) * width * 0.10
        x = left + shrink / 2
        y = top + level_h * idx
        w = width - shrink
        item = items[idx] if idx < len(items) else ""
        _add_info_box(slide, x, y, w, level_h - Inches(0.08), item, [INFO_BLUE_LIGHT, INFO_GREEN_LIGHT, INFO_AMBER_LIGHT, RGBColor(0xf8, 0xfa, 0xfc)][idx % 4], line_color=[INFO_BLUE, INFO_GREEN, INFO_AMBER, RGBColor(0xcb, 0xd5, 0xe1)][idx % 4], font_size=8.0, bold=idx == count - 1)


def _draw_horizontal_timeline(slide, items: List[str], left, top, width, height):
    count = min(max(len(items), 2), 5)
    y = top + height / 2
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.15), y, width - Inches(0.3), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = INFO_BLUE
    line.line.fill.background()
    step_w = width / count
    for idx in range(count):
        fill, line_color = INFO_COLORS[idx % len(INFO_COLORS)]
        x = left + step_w * idx + step_w / 2 - Inches(0.1)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y - Inches(0.08), Inches(0.2), Inches(0.2))
        dot.fill.solid()
        dot.fill.fore_color.rgb = line_color
        dot.line.fill.background()
        box_y = top if idx % 2 == 0 else y + Inches(0.22)
        box_h = y - top - Inches(0.16) if idx % 2 == 0 else height / 2 - Inches(0.24)
        item = items[idx] if idx < len(items) else ""
        _add_info_box(slide, left + step_w * idx + Inches(0.03), box_y, step_w - Inches(0.06), box_h, item, fill, line_color=line_color, font_size=7.6, bold=idx == 0)


def _draw_milestone_cards(slide, items: List[str], left, top, width, height):
    count = min(max(len(items), 2), 4)
    gap = Inches(0.1)
    box_w = (width - gap * (count - 1)) / count
    for idx in range(count):
        fill, line_color = INFO_COLORS[idx % len(INFO_COLORS)]
        x = left + (box_w + gap) * idx
        card = _add_info_box(slide, x, top + Inches(0.2), box_w, height - Inches(0.2), items[idx] if idx < len(items) else "", fill, line_color=line_color, font_size=7.7, bold=True)
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + box_w / 2 - Inches(0.16), top, Inches(0.32), Inches(0.32))
        badge.fill.solid()
        badge.fill.fore_color.rgb = line_color
        badge.line.fill.background()
        _set_shape_text(badge, str(idx + 1), font_size=8.5, bold=True, color=RGBColor(0xff, 0xff, 0xff))


def _draw_now_next_later(slide, items: List[str], left, top, width, height):
    labels = ["Now", "Next", "Later"]
    count = 3
    gap = Inches(0.12)
    box_w = (width - gap * (count - 1)) / count
    for idx in range(count):
        fill, line_color = INFO_COLORS[idx % 3]
        x = left + (box_w + gap) * idx
        _add_info_box(slide, x, top, box_w, Inches(0.34), labels[idx], fill, line_color=line_color, font_size=8.2, bold=True)
        text = items[idx] if idx < len(items) else ""
        _add_info_box(slide, x, top + Inches(0.44), box_w, height - Inches(0.44), text, RGBColor(0xf8, 0xfa, 0xfc), line_color=line_color, font_size=8.0)


def _draw_gantt_roadmap(slide, items: List[str], left, top, width, height):
    count = min(max(len(items), 2), 5)
    row_h = height / count
    label_w = width * 0.32
    for idx in range(count):
        fill, line_color = INFO_COLORS[idx % len(INFO_COLORS)]
        y = top + row_h * idx
        text = items[idx] if idx < len(items) else ""
        _add_info_box(slide, left, y + Inches(0.03), label_w - Inches(0.08), row_h - Inches(0.06), text[:42], RGBColor(0xf8, 0xfa, 0xfc), line_color=RGBColor(0xcb, 0xd5, 0xe1), font_size=7.4)
        bar_left = left + label_w + (width - label_w) * (idx / (count + 1)) * 0.55
        bar_w = (width - label_w) * (0.52 - idx * 0.04)
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bar_left, y + row_h * 0.25, max(bar_w, Inches(0.35)), row_h * 0.5)
        bar.fill.solid()
        bar.fill.fore_color.rgb = fill
        bar.line.color.rgb = line_color


def _draw_numbered_steps(slide, items: List[str], left, top, width, height):
    count = min(max(len(items), 2), 5)
    row_h = height / count
    for idx in range(count):
        fill, line_color = INFO_COLORS[idx % len(INFO_COLORS)]
        y = top + row_h * idx
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, y + row_h / 2 - Inches(0.15), Inches(0.3), Inches(0.3))
        badge.fill.solid()
        badge.fill.fore_color.rgb = line_color
        badge.line.fill.background()
        _set_shape_text(badge, str(idx + 1), font_size=8, bold=True, color=RGBColor(0xff, 0xff, 0xff))
        _add_info_box(slide, left + Inches(0.42), y + Inches(0.03), width - Inches(0.42), row_h - Inches(0.06), items[idx] if idx < len(items) else "", fill, line_color=line_color, font_size=8.0, bold=idx == 0)


def _draw_loop_cycle(slide, items: List[str], left, top, width, height):
    count = min(max(len(items), 3), 4)
    positions = [
        (left + width * 0.28, top),
        (left + width * 0.58, top + height * 0.32),
        (left + width * 0.28, top + height * 0.64),
        (left, top + height * 0.32),
    ]
    box_w = width * 0.42
    box_h = height * 0.28
    for idx in range(count):
        fill, line_color = INFO_COLORS[idx % len(INFO_COLORS)]
        x, y = positions[idx]
        _add_info_box(slide, x, y, box_w, box_h, items[idx] if idx < len(items) else "", fill, line_color=line_color, font_size=7.6, bold=True)
    arrows = [
        (MSO_SHAPE.RIGHT_ARROW, left + width * 0.50, top + height * 0.27),
        (MSO_SHAPE.DOWN_ARROW, left + width * 0.62, top + height * 0.50),
        (MSO_SHAPE.LEFT_ARROW, left + width * 0.46, top + height * 0.68),
        (MSO_SHAPE.UP_ARROW, left + width * 0.36, top + height * 0.50),
    ]
    for arrow_shape, arrow_x, arrow_y in arrows[:count]:
        arrow = slide.shapes.add_shape(arrow_shape, arrow_x, arrow_y, Inches(0.18), Inches(0.18))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(0x94, 0xa3, 0xb8)
        arrow.line.fill.background()


def _draw_swimlane_flow(slide, items: List[str], left, top, width, height):
    lane_count = 2
    lane_h = height / lane_count
    for lane in range(lane_count):
        fill = RGBColor(0xf8, 0xfa, 0xfc) if lane == 0 else INFO_BLUE_LIGHT
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + lane_h * lane, width, lane_h - Inches(0.03))
        band.fill.solid()
        band.fill.fore_color.rgb = fill
        band.line.color.rgb = RGBColor(0xe2, 0xe8, 0xf0)
    count = min(max(len(items), 2), 4)
    gap = Inches(0.1)
    box_w = (width - gap * (count - 1)) / count
    for idx in range(count):
        lane = idx % lane_count
        fill, line_color = INFO_COLORS[idx % len(INFO_COLORS)]
        x = left + (box_w + gap) * idx
        y = top + lane_h * lane + Inches(0.16)
        _add_info_box(slide, x, y, box_w, lane_h - Inches(0.32), items[idx] if idx < len(items) else "", fill, line_color=line_color, font_size=7.4, bold=True)


def _draw_funnel(slide, items: List[str], left, top, width, height):
    count = min(max(len(items), 3), 5)
    level_h = height / count
    for idx in range(count):
        shrink = idx * width * 0.09
        x = left + shrink / 2
        y = top + level_h * idx
        w = width - shrink
        fill, line_color = INFO_COLORS[idx % len(INFO_COLORS)]
        shape = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, x, y, w, level_h - Inches(0.04))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = line_color
        _set_shape_text(shape, items[idx] if idx < len(items) else "", font_size=7.6, bold=True)


def _draw_table_compare(slide, items: List[str], left, top, width, height):
    row_count = min(max(len(items), 2), 5)
    col_w = width / 2
    row_h = height / row_count
    for idx in range(row_count):
        fill, line_color = INFO_COLORS[idx % len(INFO_COLORS)]
        y = top + row_h * idx
        parts = re.split(r"\s*(?:vs\.?|VS|:|>|->|/)\s*", items[idx], maxsplit=1) if idx < len(items) else ["", ""]
        left_text = parts[0]
        right_text = parts[1] if len(parts) > 1 else ""
        _add_info_box(slide, left, y, col_w - Inches(0.03), row_h - Inches(0.04), left_text, fill, line_color=line_color, font_size=7.5, bold=idx == 0)
        _add_info_box(slide, left + col_w + Inches(0.03), y, col_w - Inches(0.03), row_h - Inches(0.04), right_text or items[idx], RGBColor(0xf8, 0xfa, 0xfc), line_color=RGBColor(0xcb, 0xd5, 0xe1), font_size=7.5)


def _draw_ranked_bars(slide, items: List[str], left, top, width, height):
    count = min(max(len(items), 2), 5)
    row_h = height / count
    for idx in range(count):
        fill, line_color = INFO_COLORS[idx % len(INFO_COLORS)]
        y = top + row_h * idx
        label = items[idx] if idx < len(items) else ""
        bar_w = width * (0.95 - idx * 0.13)
        _add_info_box(slide, left, y + Inches(0.03), max(bar_w, width * 0.35), row_h - Inches(0.06), f"{idx + 1}. {label}", fill, line_color=line_color, font_size=7.8, bold=idx == 0)


def _draw_before_after(slide, items: List[str], left, top, width, height):
    gap = Inches(0.16)
    col_w = (width - gap) / 2
    labels = ["Before", "After"]
    for idx in range(2):
        fill, line_color = (INFO_AMBER_LIGHT, INFO_AMBER) if idx == 0 else (INFO_GREEN_LIGHT, INFO_GREEN)
        x = left + (col_w + gap) * idx
        _add_info_box(slide, x, top, col_w, Inches(0.34), labels[idx], fill, line_color=line_color, font_size=8.3, bold=True)
        selected = items[idx::2][:3]
        box = _add_info_box(slide, x, top + Inches(0.44), col_w, height - Inches(0.44), "", RGBColor(0xf8, 0xfa, 0xfc), line_color=line_color, font_size=8.0)
        _set_list_text(box, selected, font_size=8.4)


def _draw_option_columns(slide, items: List[str], left, top, width, height):
    count = min(max(len(items), 2), 3)
    gap = Inches(0.12)
    col_w = (width - gap * (count - 1)) / count
    for idx in range(count):
        fill, line_color = INFO_COLORS[idx % len(INFO_COLORS)]
        x = left + (col_w + gap) * idx
        _add_info_box(slide, x, top, col_w, height, items[idx] if idx < len(items) else "", fill, line_color=line_color, font_size=8.0, bold=True)


def _draw_progress_bars(slide, items: List[str], left, top, width, height):
    count = min(max(len(items), 2), 5)
    row_h = height / count
    for idx in range(count):
        y = top + row_h * idx
        item = items[idx] if idx < len(items) else ""
        match = re.search(r"(\d{1,3})\s*%", item)
        pct = min(100, int(match.group(1))) if match else int(90 - idx * 12)
        _add_info_box(slide, left, y + Inches(0.02), width, row_h - Inches(0.04), item, RGBColor(0xf8, 0xfa, 0xfc), line_color=RGBColor(0xcb, 0xd5, 0xe1), font_size=7.4)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.08), y + row_h - Inches(0.12), max((width - Inches(0.16)) * pct / 100, Inches(0.08)), Inches(0.05))
        bar.fill.solid()
        bar.fill.fore_color.rgb = INFO_GREEN if pct >= 70 else INFO_AMBER
        bar.line.fill.background()


def _draw_gauge_cards(slide, items: List[str], left, top, width, height):
    count = min(max(len(items), 2), 3)
    gap = Inches(0.12)
    box_w = (width - gap * (count - 1)) / count
    for idx in range(count):
        fill, line_color = INFO_COLORS[idx % len(INFO_COLORS)]
        x = left + (box_w + gap) * idx
        _add_info_box(slide, x, top, box_w, height, "", RGBColor(0xf8, 0xfa, 0xfc), line_color=line_color)
        gauge = slide.shapes.add_shape(MSO_SHAPE.BLOCK_ARC, x + box_w * 0.22, top + Inches(0.18), box_w * 0.56, height * 0.38)
        gauge.fill.solid()
        gauge.fill.fore_color.rgb = fill
        gauge.line.color.rgb = line_color
        _add_info_box(slide, x + Inches(0.08), top + height * 0.58, box_w - Inches(0.16), height * 0.34, items[idx] if idx < len(items) else "", fill, line_color=line_color, font_size=7.4, bold=True)


def _draw_delta_callouts(slide, items: List[str], left, top, width, height):
    count = min(max(len(items), 2), 4)
    gap = Inches(0.1)
    box_w = (width - gap * (count - 1)) / count
    for idx in range(count):
        item = items[idx] if idx < len(items) else ""
        up = not re.search(r"(-|down|decrease|worse|低|減|悪)", item, flags=re.IGNORECASE)
        fill, line_color = (INFO_GREEN_LIGHT, INFO_GREEN) if up else (INFO_RED_LIGHT, INFO_RED)
        x = left + (box_w + gap) * idx
        arrow_shape = MSO_SHAPE.UP_ARROW if up else MSO_SHAPE.DOWN_ARROW
        arrow = slide.shapes.add_shape(arrow_shape, x + box_w / 2 - Inches(0.15), top, Inches(0.3), Inches(0.36))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = line_color
        arrow.line.fill.background()
        _add_info_box(slide, x, top + Inches(0.46), box_w, height - Inches(0.46), item, fill, line_color=line_color, font_size=7.6, bold=True)


def _draw_waterfall(slide, items: List[str], left, top, width, height):
    count = min(max(len(items), 3), 5)
    gap = Inches(0.08)
    bar_w = (width - gap * (count - 1)) / count
    baseline = top + height - Inches(0.12)
    for idx in range(count):
        fill, line_color = INFO_COLORS[idx % len(INFO_COLORS)]
        bar_h = height * (0.25 + 0.12 * (idx % 3))
        x = left + (bar_w + gap) * idx
        y = baseline - bar_h
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, bar_w, bar_h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = fill
        rect.line.color.rgb = line_color
        _add_info_box(slide, x, top, bar_w, y - top - Inches(0.04), items[idx] if idx < len(items) else "", RGBColor(0xf8, 0xfa, 0xfc), line_color=RGBColor(0xcb, 0xd5, 0xe1), font_size=7.0)


def _draw_risk_register(slide, items: List[str], left, top, width, height):
    headers = ["Risk", "Impact", "Action"]
    header_h = Inches(0.28)
    col_w = width / 3
    for idx, header in enumerate(headers):
        _add_info_box(slide, left + col_w * idx, top, col_w - Inches(0.03), header_h, header, [INFO_RED_LIGHT, INFO_AMBER_LIGHT, INFO_GREEN_LIGHT][idx], line_color=[INFO_RED, INFO_AMBER, INFO_GREEN][idx], font_size=7.5, bold=True)
    row_count = min(max(len(items), 2), 4)
    row_h = (height - header_h - Inches(0.06)) / row_count
    for row in range(row_count):
        parts = [part.strip() for part in re.split(r"\s*(?:\||/|:|->)\s*", items[row] if row < len(items) else "", maxsplit=2)]
        for col in range(3):
            text = parts[col] if col < len(parts) else ""
            _add_info_box(slide, left + col_w * col, top + header_h + Inches(0.06) + row_h * row, col_w - Inches(0.03), row_h - Inches(0.03), text, RGBColor(0xf8, 0xfa, 0xfc), line_color=RGBColor(0xcb, 0xd5, 0xe1), font_size=6.8)


def _draw_heatmap(slide, items: List[str], left, top, width, height):
    rows = cols = 3
    gap = Inches(0.04)
    cell_w = (width - gap * (cols - 1)) / cols
    cell_h = (height - gap * (rows - 1)) / rows
    heat = [INFO_GREEN_LIGHT, INFO_AMBER_LIGHT, INFO_RED_LIGHT]
    lines = [INFO_GREEN, INFO_AMBER, INFO_RED]
    for row in range(rows):
        for col in range(cols):
            idx = row * cols + col
            level = min(2, max(0, row + col - 1))
            text = items[idx] if idx < len(items) else ""
            _add_info_box(slide, left + (cell_w + gap) * col, top + (cell_h + gap) * row, cell_w, cell_h, text, heat[level], line_color=lines[level], font_size=6.8, bold=bool(text))


def _draw_escalation_ladder(slide, items: List[str], left, top, width, height):
    count = min(max(len(items), 3), 5)
    step_h = height / count
    for idx in range(count):
        fill, line_color = [INFO_GREEN_LIGHT, INFO_BLUE_LIGHT, INFO_AMBER_LIGHT, INFO_RED_LIGHT, INFO_PURPLE_LIGHT][idx % 5], [INFO_GREEN, INFO_BLUE, INFO_AMBER, INFO_RED, INFO_PURPLE][idx % 5]
        x = left + width * idx * 0.08
        y = top + height - step_h * (idx + 1)
        w = width - width * idx * 0.08
        _add_info_box(slide, x, y, w, step_h - Inches(0.06), items[idx] if idx < len(items) else "", fill, line_color=line_color, font_size=7.5, bold=idx == count - 1)


def _draw_priority_quadrants(slide, items: List[str], left, top, width, height):
    _draw_matrix(slide, items, left, top, width, height)
    _add_info_box(slide, left + Inches(0.06), top + Inches(0.04), Inches(0.74), Inches(0.24), "High", INFO_RED_LIGHT, line_color=INFO_RED, font_size=6.5, bold=True)
    _add_info_box(slide, left + width - Inches(0.82), top + height - Inches(0.3), Inches(0.74), Inches(0.24), "Low", INFO_GREEN_LIGHT, line_color=INFO_GREEN, font_size=6.5, bold=True)


def _draw_decision_matrix(slide, items: List[str], left, top, width, height):
    headers = ["Option", "Criteria", "Decision"]
    header_h = Inches(0.3)
    col_w = width / 3
    for idx, header in enumerate(headers):
        fill, line_color = INFO_COLORS[idx]
        _add_info_box(slide, left + col_w * idx, top, col_w - Inches(0.03), header_h, header, fill, line_color=line_color, font_size=7.2, bold=True)
    rows = min(max(len(items), 2), 4)
    row_h = (height - header_h - Inches(0.06)) / rows
    for row in range(rows):
        text = items[row] if row < len(items) else ""
        parts = [part.strip() for part in re.split(r"\s*(?:\||/|:|->)\s*", text, maxsplit=2)]
        for col in range(3):
            _add_info_box(slide, left + col_w * col, top + header_h + Inches(0.06) + row_h * row, col_w - Inches(0.03), row_h - Inches(0.03), parts[col] if col < len(parts) else "", RGBColor(0xf8, 0xfa, 0xfc), line_color=RGBColor(0xcb, 0xd5, 0xe1), font_size=6.8)


def _draw_hub_spoke(slide, items: List[str], left, top, width, height):
    center_w = min(width * 0.38, Inches(1.35))
    center_h = min(height * 0.36, Inches(0.75))
    cx = left + width / 2 - center_w / 2
    cy = top + height / 2 - center_h / 2
    center = items[0] if items else ""
    spokes = items[1:7] or items[:6]
    positions = [
        (left, top),
        (left + width - center_w, top),
        (left, top + height - center_h),
        (left + width - center_w, top + height - center_h),
        (left + width / 2 - center_w / 2, top),
        (left + width / 2 - center_w / 2, top + height - center_h),
    ]
    center_x = cx + center_w / 2
    center_y = cy + center_h / 2
    line_specs = [
        (left + center_w, center_y, cx - (left + center_w), Inches(0.02)),
        (cx + center_w, center_y, left + width - center_w - (cx + center_w), Inches(0.02)),
        (center_x, top + center_h, Inches(0.02), cy - (top + center_h)),
        (center_x, cy + center_h, Inches(0.02), top + height - center_h - (cy + center_h)),
    ]
    for x, y, w, h in line_specs:
        if w <= 0 or h <= 0:
            continue
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0xcb, 0xd5, 0xe1)
        line.line.fill.background()
    _add_info_box(slide, cx, cy, center_w, center_h, center, INFO_BLUE_LIGHT, line_color=INFO_BLUE, font_size=8.0, bold=True)
    for idx, item in enumerate(spokes[:6]):
        fill, line_color = INFO_COLORS[(idx + 1) % len(INFO_COLORS)]
        x, y = positions[idx]
        _add_info_box(slide, x, y, center_w, center_h, item, fill, line_color=line_color, font_size=7.0, bold=False)


def _resolve_visual_variant(style: str, requested: str, items: List[str], width, height) -> str:
    requested = (requested or "auto").strip()
    if requested != "auto" and requested in VISUAL_VARIANTS_BY_STYLE.get(style, ()):
        return requested
    if style == "timeline":
        if any(re.search(r"\d+%|q[1-4]|\d{4}[/-]\d{1,2}", item, flags=re.IGNORECASE) for item in items):
            return "gantt_roadmap"
        if len(items) == 3:
            return "now_next_later"
        if width > height * 2.0:
            return "horizontal_timeline"
        return "phase_bands" if len(items) <= 3 else "vertical_timeline"
    if style == "process":
        if len(items) >= 5:
            return "funnel"
        return "chevron_flow" if width > height * 1.8 else "vertical_flow"
    if style == "comparison":
        if len(items) == 2:
            return "before_after"
        return "pros_cons" if len(items) >= 6 else "cards_2x2"
    if style == "kpi":
        numberish = sum(1 for item in items if re.search(r"\d", item))
        if any("%" in item for item in items):
            return "progress_bars"
        return "big_numbers" if numberish >= 2 and len(items) <= 3 else "metric_cards"
    if style == "risk":
        return "heatmap" if len(items) >= 6 else "cause_impact_mitigation"
    if style == "matrix":
        return "priority_quadrants"
    if style == "summary":
        return "hub_spoke" if len(items) >= 5 else ("pyramid" if len(items) <= 3 else "bands")
    return "bands"


def _enhance_template_slide_infographic(slide, slide_data: SlideNode):
    if any(content.image_prompt or content.use_user_image for content in slide_data.placeholders):
        return
    style = _infer_visual_style(slide_data)
    if style == "none":
        return
    target = _visual_placeholder_candidate(slide, slide_data)
    if target is None:
        return
    target_idx = target.placeholder_format.idx
    source_text = _placeholder_text(slide_data, {target_idx}) or _placeholder_text(
        slide_data,
        {content.idx for content in slide_data.placeholders},
    )
    items = _split_visual_items(source_text)
    if len(items) < 2:
        return
    left, top, width, height = target.left, target.top, target.width, target.height
    _delete_shape(target)
    variant = _resolve_visual_variant(style, getattr(slide_data, "visual_variant", "auto"), items, width, height)
    if style == "timeline":
        if variant == "horizontal_timeline":
            _draw_horizontal_timeline(slide, items, left, top, width, height)
        elif variant == "milestone_cards":
            _draw_milestone_cards(slide, items, left, top, width, height)
        elif variant == "now_next_later":
            _draw_now_next_later(slide, items, left, top, width, height)
        elif variant == "gantt_roadmap":
            _draw_gantt_roadmap(slide, items, left, top, width, height)
        elif variant == "phase_bands":
            _draw_phase_bands(slide, items, left, top, width, height)
        else:
            _draw_timeline(slide, items, left, top, width, height)
    elif style == "process":
        if variant == "numbered_steps":
            _draw_numbered_steps(slide, items, left, top, width, height)
        elif variant == "loop_cycle":
            _draw_loop_cycle(slide, items, left, top, width, height)
        elif variant == "swimlane_flow":
            _draw_swimlane_flow(slide, items, left, top, width, height)
        elif variant == "funnel":
            _draw_funnel(slide, items, left, top, width, height)
        elif variant == "chevron_flow":
            _draw_chevron_flow(slide, items, left, top, width, height)
        else:
            _draw_flow(slide, items, left, top, width, height)
    elif style == "comparison":
        if variant == "table_compare":
            _draw_table_compare(slide, items, left, top, width, height)
        elif variant == "ranked_bars":
            _draw_ranked_bars(slide, items, left, top, width, height)
        elif variant == "before_after":
            _draw_before_after(slide, items, left, top, width, height)
        elif variant == "option_columns":
            _draw_option_columns(slide, items, left, top, width, height)
        elif variant == "pros_cons":
            _draw_pros_cons(slide, items, left, top, width, height)
        elif variant == "scorecard":
            _draw_scorecard(slide, items, left, top, width, height)
        else:
            _draw_kpi_grid(slide, items, style, left, top, width, height)
    elif style == "kpi":
        if variant == "progress_bars":
            _draw_progress_bars(slide, items, left, top, width, height)
        elif variant == "gauge_cards":
            _draw_gauge_cards(slide, items, left, top, width, height)
        elif variant == "delta_callouts":
            _draw_delta_callouts(slide, items, left, top, width, height)
        elif variant == "waterfall":
            _draw_waterfall(slide, items, left, top, width, height)
        elif variant == "big_numbers":
            _draw_big_numbers(slide, items, left, top, width, height)
        elif variant == "scorecard":
            _draw_scorecard(slide, items, left, top, width, height)
        else:
            _draw_kpi_grid(slide, items, style, left, top, width, height)
    elif style == "risk":
        if variant == "risk_register":
            _draw_risk_register(slide, items, left, top, width, height)
        elif variant == "heatmap":
            _draw_heatmap(slide, items, left, top, width, height)
        elif variant == "escalation_ladder":
            _draw_escalation_ladder(slide, items, left, top, width, height)
        elif variant == "cause_impact_mitigation":
            _draw_cause_impact_mitigation(slide, items, left, top, width, height)
        else:
            _draw_matrix(slide, items, left, top, width, height)
    elif style == "matrix":
        if variant == "priority_quadrants":
            _draw_priority_quadrants(slide, items, left, top, width, height)
        elif variant == "decision_matrix":
            _draw_decision_matrix(slide, items, left, top, width, height)
        else:
            _draw_matrix(slide, items, left, top, width, height)
    else:
        if variant == "hub_spoke":
            _draw_hub_spoke(slide, items, left, top, width, height)
        elif variant == "pyramid":
            _draw_pyramid(slide, items, left, top, width, height)
        else:
            _draw_summary_bands(slide, items, left, top, width, height)


def render_pptx_slide(prs: Presentation, slide_data: SlideNode, font_size_offset: int = 0, current_index: int = 0, total_slides: int = 0, image_paths: Dict[int, str] = None, has_template: bool = False, blank_layout = None):
    """SlideNodeデータから テンプレートのプレースホルダーへ直接データを流し込みレンダリングする"""
    if blank_layout is None:
        blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    # テンプレートを使用している場合
    if has_template:
        ph_map = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
        # 共通要素：スライド番号の書き込み (SLIDE_NUMBER タイプ)
        for ph in slide.placeholders:
            if ph.placeholder_format.type == 13:  # SLIDE_NUMBER
                ph.text = f"{current_index + 1} / {total_slides}"
                for p in ph.text_frame.paragraphs:
                    p.font.name = "Meiryo"
                    p.font.size = Pt(10)
        # LLMから出力されたプレースホルダー情報を処理
        for content in slide_data.placeholders:
            idx = content.idx
            if idx not in ph_map:
                continue
            ph = ph_map[idx]
            # テキストコンテンツの流し込み
            if content.text_content:
                ph.text = content.text_content
                _format_text_placeholder(ph, font_size_offset=font_size_offset)
            # 画像コンテンツの流し込み
            elif image_paths and idx in image_paths:
                img_path = image_paths[idx]
                if os.path.exists(img_path):
                    try:
                        # プレースホルダーの物理的な位置とサイズを取得
                        left = ph.left
                        top = ph.top
                        width = ph.width
                        height = ph.height
                        # プレースホルダーの位置に画像を挿入
                        slide.shapes.add_picture(img_path, left, top, width, height)
                        # 元のプレースホルダー（ゴースト）を削除
                        sp = ph._element
                        sp.getparent().remove(sp)
                    except Exception as e:
                        state_manager.add_debug_log(f"[PPTXAgent] Failed to insert image to placeholder idx={idx}: {e}", "warning")
        # AIが使用しなかった不要なプレースホルダー（「タイトルを追加」や「画像を追加」の枠線）を完全に消去する
        used_idxs = set(content.idx for content in slide_data.placeholders)
        content_ph_types = {1, 2, 3, 4, 7, 12, 18} # 主要なコンテンツ用プレースホルダータイプ
        for shape in list(slide.shapes):
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                idx = shape.placeholder_format.idx
                if ph_type in content_ph_types and idx not in used_idxs:
                    try:
                        sp = shape._element
                        sp.getparent().remove(sp)
                    except Exception as e:
                        pass
        _enhance_template_slide_infographic(slide, slide_data)
    else:
        # フォールバック: テンプレートがない場合のマニュアル簡易レンダリング
        from pptx.dml.color import RGBColor
        COLOR_BG = RGBColor(0xff, 0xff, 0xff)
        COLOR_TEXT_MAIN = RGBColor(0x0f, 0x17, 0x2a)
        bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = COLOR_BG
        bg_shape.line.fill.background()
        tx_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.8))
        tx_title.text_frame.text = slide_data.title
        tx_title.text_frame.paragraphs[0].font.name = "Meiryo"
        tx_title.text_frame.paragraphs[0].font.size = Pt(32)
        tx_title.text_frame.paragraphs[0].font.bold = True
        tx_title.text_frame.paragraphs[0].font.color.rgb = COLOR_TEXT_MAIN
        tx_body = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.5))
        tf = tx_body.text_frame
        tf.word_wrap = True
        for content in slide_data.placeholders:
            if content.text_content:
                p = tf.add_paragraph()
                p.text = content.text_content
                p.font.name = "Meiryo"
                p.font.size = Pt(18 + font_size_offset)

# --- Playwright 同期バリデーション実行ヘルパー ---


def validate_slide_overflow(temp_dir: str, slides: List[SlideNode], offsets: List[int]) -> List[dict]:
    results = []
    browser = None
    try:
        with sync_playwright() as p:
            launch_attempts = [
                {},                      # 1. Playwright デフォルト Chromium
                {"channel": "msedge"},   # 2. システムインストール済み Edge
                {"channel": "chrome"}    # 3. システムインストール済み Chrome
            ]
            for attempt in launch_attempts:
                try:
                    browser = p.chromium.launch(headless=True, **attempt)
                    state_manager.add_debug_log(f"[PPTXAgent] Launched browser successfully with settings: {attempt}")
                    break
                except Exception as e:
                    state_manager.add_debug_log(f"[PPTXAgent] Browser launch attempt failed ({attempt}): {e}", "warning")
            if browser:
                for slide in slides:
                    html_path = os.path.join(temp_dir, f"{slide.slide_number:02d}.html")
                    res = validate_single_slide(browser, html_path, slide.slide_number)
                    results.append(res)
                browser.close()
            else:
                state_manager.add_debug_log(
                    "[PPTXAgent] All browser launch attempts failed. Skipping geometry validation and proceeding with default offsets.", 
                    "warning"
                )
                for slide in slides:
                    results.append({"overflowed": False, "slide_number": slide.slide_number})
    except Exception as e:
        state_manager.add_debug_log(f"[PPTXAgent] Playwright validation encountered unexpected error: {e}", "warning")
        # 全滅時の緊急フォールバック結果の構築
        if not results:
            for slide in slides:
                results.append({"overflowed": False, "slide_number": slide.slide_number})
    return results


def validate_single_slide(browser, html_path: str, slide_number: int) -> dict:
    page = browser.new_page()
    page.set_viewport_size({"width": 1280, "height": 720})
    file_uri = Path(os.path.abspath(html_path)).as_uri()
    page.goto(file_uri)
    page.wait_for_load_state("networkidle")
    overflow_js = """
    () => {
        const targets = document.querySelectorAll('.dynamic-text-block');
        for (let el of targets) {
            if (el.scrollHeight > el.clientHeight || el.scrollWidth > el.clientWidth) {
                return {
                    overflowed: true,
                    element_id: el.id,
                    text_content: el.innerText,
                    available_height: el.clientHeight,
                    required_height: el.scrollHeight
                };
            }
        }
        return { overflowed: false };
    }
    """
    res = page.evaluate(overflow_js)
    res["slide_number"] = slide_number
    page.close()
    return res


def _set_textbox_text(shape, text: str, font_size: float, color: RGBColor, bold: bool = False, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = "Meiryo"
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color


def _add_reference_row(slide, y, reference: ReferenceEntry, fill_color: RGBColor, line_color: RGBColor):
    text_color = RGBColor(0x1f, 0x29, 0x37)
    muted_color = RGBColor(0x64, 0x74, 0x8b)
    green_color = RGBColor(0x0f, 0x76, 0x66)
    row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.58), y, Inches(12.15), Inches(0.66))
    row_bg.fill.solid()
    row_bg.fill.fore_color.rgb = fill_color
    row_bg.line.color.rgb = line_color
    row_bg.line.width = Pt(0.5)
    ref_id = reference.reference_id or "-"
    source_type = (reference.source_type or "other").upper()
    used = ", ".join(f"S{num}" for num in reference.used_slides[:6]) if reference.used_slides else "-"
    title = _shorten_text(reference.title, 86)
    snippet = _shorten_text(reference.snippet, 150)
    uri = _shorten_text(reference.uri, 110) if reference.uri else ""
    id_box = slide.shapes.add_textbox(Inches(0.7), y + Inches(0.08), Inches(0.58), Inches(0.48))
    _set_textbox_text(id_box, ref_id, 8.4, green_color, bold=True, align=PP_ALIGN.CENTER)
    type_box = slide.shapes.add_textbox(Inches(1.28), y + Inches(0.08), Inches(0.95), Inches(0.48))
    _set_textbox_text(type_box, source_type, 6.6, muted_color, bold=True, align=PP_ALIGN.CENTER)
    body_box = slide.shapes.add_textbox(Inches(2.3), y + Inches(0.06), Inches(7.95), Inches(0.54))
    body_text = f"{title}\n{snippet}" + (f"\n{uri}" if uri else "")
    _set_textbox_text(body_box, body_text, 6.6, text_color)
    used_box = slide.shapes.add_textbox(Inches(10.35), y + Inches(0.08), Inches(2.1), Inches(0.48))
    _set_textbox_text(used_box, used, 7.0, muted_color, align=PP_ALIGN.CENTER)


def add_references_slide(prs: Presentation, layout, references: List[ReferenceEntry]) -> bool:
    visible_references = [
        ref for ref in references
        if ref.title or ref.uri or ref.snippet
    ][:8]
    if not visible_references:
        return False
    slide = prs.slides.add_slide(layout)
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            _delete_shape(shape)
    primary = RGBColor(0x1e, 0x3a, 0x8a)
    muted = RGBColor(0x64, 0x74, 0x8b)
    line = RGBColor(0xd7, 0xde, 0xe8)
    fill_even = RGBColor(0xf8, 0xfa, 0xfc)
    fill_odd = RGBColor(0xff, 0xff, 0xff)
    title_box = slide.shapes.add_textbox(Inches(0.68), Inches(0.32), Inches(7.5), Inches(0.42))
    _set_textbox_text(title_box, "出典・参照情報", 20, primary, bold=True)
    note_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.78), Inches(9.5), Inches(0.25))
    _set_textbox_text(note_box, "主要なWeb出典・添付資料のみを掲載。詳細は生成フォルダの source_brief / references JSON を参照。", 7.2, muted)
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.58), Inches(1.12), Inches(12.15), Inches(0.28))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(0xeb, 0xf1, 0xf7)
    header.line.color.rgb = line
    header.line.width = Pt(0.5)
    for x, width, label, align in [
        (0.7, 0.58, "ID", PP_ALIGN.CENTER),
        (1.28, 0.95, "種別", PP_ALIGN.CENTER),
        (2.3, 7.95, "出典 / 根拠メモ", PP_ALIGN.LEFT),
        (10.35, 2.1, "使用スライド", PP_ALIGN.CENTER),
    ]:
        box = slide.shapes.add_textbox(Inches(x), Inches(1.16), Inches(width), Inches(0.2))
        _set_textbox_text(box, label, 6.5, muted, bold=True, align=align)
    y = Inches(1.48)
    for index, reference in enumerate(visible_references):
        _add_reference_row(
            slide=slide,
            y=y + Inches(0.7 * index),
            reference=reference,
            fill_color=fill_even if index % 2 == 0 else fill_odd,
            line_color=line,
        )
    if len(references) > len(visible_references):
        more_box = slide.shapes.add_textbox(Inches(0.7), Inches(7.02), Inches(12), Inches(0.22))
        _set_textbox_text(more_box, f"ほか {len(references) - len(visible_references)} 件の参照情報は *_references.json に保存されています。", 7.0, muted)
    return True


def reset_pptx_open_view(pptx_path: str) -> None:
    """Force generated PPTX files to open in normal slide editing view."""
    import tempfile
    import zipfile
    view_props_name = "ppt/viewProps.xml"
    target_dir = os.path.dirname(os.path.abspath(pptx_path)) or None
    fd, temp_path = tempfile.mkstemp(suffix=".pptx", dir=target_dir)
    os.close(fd)
    try:
        changed = False
        with zipfile.ZipFile(pptx_path, "r") as zin, zipfile.ZipFile(temp_path, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename == view_props_name:
                    text = data.decode("utf-8", errors="replace")
                    if "lastView=" in text:
                        text = re.sub(r'lastView="[^"]+"', 'lastView="sldThumbnailView"', text, count=1)
                    else:
                        text = text.replace("<p:viewPr ", '<p:viewPr lastView="sldThumbnailView" ', 1)
                    data = text.encode("utf-8")
                    changed = True
                zout.writestr(item, data)
        if changed:
            os.replace(temp_path, pptx_path)
            state_manager.add_debug_log("[PPTXAgent] Reset PPTX open view to normal slide view.")
        else:
            try:
                os.remove(temp_path)
            except Exception:
                pass
    except Exception as e:
        try:
            os.remove(temp_path)
        except Exception:
            pass
        state_manager.add_debug_log(f"[PPTXAgent] Failed to reset PPTX open view: {e}", "warning")

# --- PPTXAgent 本体クラス ---


class PPTXAgent:
    def __init__(self, client):
        self.llm_clients = llm_router.coerce_llm_clients(client)
        state_manager.add_debug_log("[PPTXAgent] Initialized successfully.")

    def _run_research_context(
        self,
        chat_contents: List[Any],
        tools_config: Optional[List[Any]],
        model_id: str,
    ) -> tuple[str, Optional[dict]]:
        if not tools_config:
            return "", None
        try:
            research_contents = [_clone_content_for_pptx(content) for content in chat_contents]
            research_contents.append(
                types.Content(
                    role="user",
                    parts=[
                        types.Part.from_text(
                            text=(
                                "PowerPoint資料化の前提調査として、会話の目的に必要な最新情報・事実関係・"
                                "重要な出典候補を簡潔に整理してください。"
                                "添付資料や会話内容と矛盾しない範囲で、スライドに入れるべき事実だけを抽出してください。"
                            )
                        )
                    ],
                )
            )
            research_config = types.GenerateContentConfig(
                system_instruction=(
                    "あなたは資料作成前のリサーチ担当です。最新性が必要な論点は検索し、"
                    "スライドに使える事実、日付、比較軸、注意点を日本語で簡潔に整理してください。"
                    "出典候補がわかる場合はURLや媒体名も残してください。"
                ),
                max_output_tokens=8192,
                temperature=0.1,
                tools=tools_config,
            )
            research_result = llm_router.generate_content_with_route(
                llm_clients=self.llm_clients,
                model_id=model_id,
                contents=research_contents,
                config=research_config,
                mode="pptx_research",
                logger=state_manager.add_debug_log,
            )
            grounding_text = _format_grounding_metadata(research_result.grounding_metadata)
            if research_result.text or grounding_text:
                state_manager.add_debug_log("[PPTXAgent] Added web/tool research context for PPTX source brief.")
                return (
                    (research_result.text or "")[:12000]
                    + ("\n\n【検索メタデータ】\n" + grounding_text if grounding_text else "")
                ), research_result.grounding_metadata
        except Exception as e:
            state_manager.add_debug_log(f"[PPTXAgent] PPTX research step skipped: {e}", "warning")
        return "", None

    def _build_source_brief(
        self,
        chat_contents: List[Any],
        conversation_excerpt: str,
        attachment_summary: str,
        file_attachments_meta: Optional[List[dict]],
        research_context: str,
        grounding_metadata: Optional[dict],
        materialized_system_instruction: str,
        model_id: str,
    ) -> PresentationSourceBrief:
        source_system = ""
        if materialized_system_instruction:
            source_system = (
                "\n\n【元チャットのシステム/モード指示】\n"
                f"{materialized_system_instruction[:6000]}"
            )
        brief_contents = [_clone_content_for_pptx(content) for content in chat_contents]
        brief_contents.append(
            types.Content(
                role="user",
                parts=[
                    types.Part.from_text(
                        text=(
                            "これまでの会話、添付ファイル、添付画像、検索結果、キャンバス情報を統合し、"
                            "PowerPoint資料を作るための source brief JSON を作成してください。\n"
                            "重要: テンプレートの初期文言を材料として扱わず、会話・添付・検索で確認できる情報だけを材料にしてください。\n"
                            "重要: evidence_notes には、実際に与えられた添付名、会話内の明示情報、検索メタデータに含まれるURL/媒体名だけを書いてください。"
                            "確認できない公式ブログ名や出典名を推測で書いてはいけません。\n"
                            "重要: source_inventory には、利用可能な材料を種類別に漏れなく列挙してください。"
                            "検索結果がある場合は『Web検索結果』、添付がある場合は添付ファイル名、画像がある場合は画像名を必ず含めてください。\n"
                            "重要: すべての材料をそのまま詰め込むのではなく、資料目的に必要な根拠・数値・日付・比較軸を台帳化してください。\n"
                            "重要: 添付画像や生成画像は、内容理解に役立つ場合だけ visual_assets / image_policy に残してください。\n"
                            "重要: source_coverage_units には、最終PPTXに反映可否を追跡すべき原子情報を入れてください。"
                            "会話の明示依頼、添付ファイル/画像から読み取れる事実、検索で得た主要事実、制約条件、未確定だが注意すべき点を短い単位で列挙します。\n"
                            "重要: source_coverage_units は短すぎる見出しではなく、スライド本文に転記できる程度の具体性を持たせてください。"
                            "十分な材料がある場合は12〜30件程度に分解し、数値・日付・仕組み・比較軸・例外条件を落とさないでください。\n"
                            "重要: references には、最終PPTXの出典スライドに載せる参照情報を作成してください。"
                            "Web検索結果は title / uri / snippet を必ず入れ、snippet にはその出典から確認できる本文スニペットまたは短い根拠要約を書いてください。"
                            "添付資料はファイル名を title に入れ、snippet に資料内容や利用した根拠の短い要約を書いてください。"
                            "reference_id は R1, R2 のように連番にしてください。used_slides は後段で補完されるため空でも構いません。\n"
                            + (
                                "\n\n【会話ログ抜粋】\n"
                                f"{conversation_excerpt}"
                                if conversation_excerpt
                                else ""
                            )
                            + attachment_summary
                            + (
                                "\n\n【PPTX生成前リサーチ結果】\n"
                                f"{research_context}"
                                if research_context
                                else ""
                            )
                            + (
                                "\n\n【構造化検索メタデータ】\n"
                                f"{json.dumps(grounding_metadata, ensure_ascii=False, indent=2)}"
                                if grounding_metadata
                                else ""
                            )
                            + source_system
                        )
                    )
                ],
            )
        )
        brief_config = types.GenerateContentConfig(
            system_instruction=(
                "あなたは資料作成前の編集長です。会話、添付、画像、検索結果を統合し、"
                "スライド構成に使える材料台帳を作ります。"
                "事実と推測を混ぜず、根拠が弱いものは gaps_or_uncertainties に分離してください。"
                "最終資料で必ず反映すべき観点は coverage_requirements に入れてください。"
                "source_coverage_units には、coverage_requirements より細かい情報単位を作り、後続のスライド構成で coverage_refs として追跡できるようにしてください。"
                "source_coverage_units は監査用の短いIDだけでなく、本文化できる具体文にしてください。情報が多い場合は過度に統合せず、別々の単位に分けてください。"
                "references は配布用PPTXに載せる出典一覧です。Web URL、添付ファイル名、会話由来、Canvas由来を区別し、根拠スニペットを短く残してください。"
            ),
            max_output_tokens=32768,
            temperature=0.1,
            response_mime_type="application/json",
            response_schema=PresentationSourceBrief,
        )
        result = llm_router.generate_content_with_route(
            llm_clients=self.llm_clients,
            model_id=model_id,
            contents=brief_contents,
            config=brief_config,
            mode="pptx_source_brief",
            logger=state_manager.add_debug_log,
        )
        brief = PresentationSourceBrief.model_validate_json(result.text)
        state_manager.add_debug_log(
            f"[PPTXAgent] Source brief built: facts={len(brief.key_facts)}, "
            f"evidence={len(brief.evidence_notes)}, coverage={len(brief.coverage_requirements)}, "
            f"units={len(brief.source_coverage_units)}."
        )
        if not brief.source_coverage_units:
            fallback_units = []
            fallback_units.extend(brief.coverage_requirements[:20])
            fallback_units.extend(brief.key_facts[:30])
            fallback_units.extend(brief.visual_assets[:10])
            brief.source_coverage_units = list(dict.fromkeys(unit for unit in fallback_units if unit))
            state_manager.add_debug_log(
                f"[PPTXAgent] Filled source coverage units from brief fallback: {len(brief.source_coverage_units)}."
            )
        brief.references = _finalize_reference_entries(
            references=brief.references,
            file_attachments_meta=file_attachments_meta,
            grounding_metadata=grounding_metadata,
        )
        required_sources = ["会話履歴（ユーザー依頼とこれまでの回答）"]
        if attachment_summary:
            required_sources.append("添付ファイル/画像情報")
        if research_context:
            required_sources.append("Web検索/grounding結果")
        existing_sources = set(brief.source_inventory)
        for source_name in required_sources:
            if source_name not in existing_sources:
                brief.source_inventory.append(source_name)
        return brief

    def _generate_presentation_structure(
        self,
        brief: PresentationSourceBrief,
        system_instruction: str,
        gen_config: types.GenerateContentConfig,
        model_id: str,
        revision_instruction: str = "",
    ) -> PresentationDSLSchema:
        min_body_slides = _minimum_body_slide_count(brief)
        prompt = (
            "以下の source brief を唯一の材料台帳として、最高品質のスライド資料構成案JSONを出力してください。\n"
            f"本文スライドは最低 {min_body_slides} 枚作成してください。表紙はテンプレート側で別途作成されるため、このJSONには本文・比較・結論スライドを入れてください。\n"
            "recommended_storyline がある場合は、表紙を除く各章を原則として別スライドにしてください。複数章を1枚に統合してはいけません。\n"
            "1枚のスライドで扱う source_coverage_units は原則2件までにしてください。比較・結論スライドだけは例外として複数件を横断して構いません。\n"
            "必ず coverage_requirements を本文スライドへ反映してください。\n"
            "source_coverage_units も可能な限り本文スライドへ反映してください。\n"
            "各スライドの coverage_refs には、そのスライドで扱った coverage_requirements / source_coverage_units / key_facts の文言を短く引用してください。\n"
            "references に含まれる出典を使う場合は、coverage_refs に R1, R2 などの reference_id も含めてください。\n"
            "coverage_refs は監査用メタデータであり、本文の代替ではありません。重要情報は必ず placeholders の text_content または image_prompt に可視情報として入れてください。\n"
            "各本文スライドはタイトルだけ、短いラベルだけ、空の区切りスライドだけにしてはいけません。\n"
            "情報量が多い場合はスライド数を増やして構いません。重要情報を削って見栄えだけを優先してはいけません。\n"
            "key_facts と evidence_notes にない断定は避けてください。\n"
            "evidence_notes が弱い項目は断定調ではなく『要確認』『現時点の情報』として扱ってください。\n"
            "source_inventory が少ない場合は、薄い資料を無理に水増しせず、不確実点と次アクションを明示してください。\n"
            "gaps_or_uncertainties は、必要なら注意点として短く扱ってください。\n"
            "テンプレートの初期文言を出力してはいけません。\n\n"
            f"【source brief】\n{_brief_to_text(brief)}"
        )
        if revision_instruction:
            prompt += f"\n\n【修正指示】\n{revision_instruction}"
        result = llm_router.generate_content_with_route(
            llm_clients=self.llm_clients,
            model_id=model_id,
            contents=[types.Content(role="user", parts=[types.Part.from_text(text=prompt)])],
            config=gen_config,
            mode="pptx_structure",
            logger=state_manager.add_debug_log,
        )
        presentation_data = PresentationDSLSchema.model_validate_json(result.text)
        state_manager.add_debug_log(
            f"[PPTXAgent] Storyline established. Title: {presentation_data.presentation_title}, "
            f"Slides: {len(presentation_data.slides)}"
        )
        return presentation_data

    def _audit_content_coverage(
        self,
        brief: PresentationSourceBrief,
        presentation_data: PresentationDSLSchema,
        model_id: str,
    ) -> PresentationCoverageAudit:
        required_items = _required_coverage_items(brief)
        if not required_items:
            return PresentationCoverageAudit(status="pass")
        visible_text = _deck_visible_text(presentation_data.slides)
        min_body_slides = _minimum_body_slide_count(brief)
        trace_map = [
            {
                "slide_number": slide.slide_number,
                "title": slide.title,
                "coverage_refs": slide.coverage_refs,
            }
            for slide in presentation_data.slides
        ]
        low_content_slides = _low_content_slide_numbers(presentation_data.slides)
        prompt = (
            "You are auditing content coverage for a generated PowerPoint structure.\n"
            "The deck must use all available relevant information, not only produce attractive slides.\n"
            "Audit whether each required item is clearly represented in visible slide titles, placeholder text, or image prompts.\n"
            "Important: coverage_refs are traceability metadata only. Do not count coverage_refs as visible content.\n"
            "Mark an item missing if it is absent from visible content or only appears in coverage_refs.\n"
            "Mark it weakly covered if it appears but loses important qualifiers, numbers, dates, sources, mechanisms, tradeoffs, or user constraints.\n"
            "Mark it overcompressed if several distinct source facts were collapsed into a generic statement that loses useful information.\n"
            "Fail the audit if any slide is empty or has almost no body content unless it is clearly a cover/back-cover slide.\n"
            "Return strict JSON following the schema. If anything material is missing, status must be fail.\n\n"
            f"Required items:\n{json.dumps(required_items, ensure_ascii=False, indent=2)}\n\n"
            f"Source brief:\n{_brief_to_text(brief)}\n\n"
            f"Minimum body slide count required by deterministic check: {min_body_slides}\n"
            f"Actual generated body slide count: {len(presentation_data.slides)}\n\n"
            f"Visible deck text:\n{visible_text}\n\n"
            f"Traceability metadata only, not visible content:\n{json.dumps(trace_map, ensure_ascii=False, indent=2)}\n\n"
            f"Slides with low visible body content by deterministic check: {low_content_slides}\n\n"
            f"Generated structure:\n{json.dumps(presentation_data.model_dump(), ensure_ascii=False, indent=2)}"
        )
        audit_config = types.GenerateContentConfig(
            max_output_tokens=16384,
            temperature=0.0,
            response_mime_type="application/json",
            response_schema=PresentationCoverageAudit,
        )
        try:
            result = llm_router.generate_content_with_route(
                llm_clients=self.llm_clients,
                model_id=model_id,
                contents=[types.Content(role="user", parts=[types.Part.from_text(text=prompt)])],
                config=audit_config,
                mode="pptx_content_coverage_audit",
                logger=state_manager.add_debug_log,
            )
            audit = PresentationCoverageAudit.model_validate_json(result.text)
        except Exception as e:
            state_manager.add_debug_log(f"[PPTXAgent] Content coverage audit skipped: {e}", "warning")
            return PresentationCoverageAudit(
                status="warning",
                recommendations=["Content coverage audit could not be completed."],
                revision_instruction="Preserve and reflect all coverage_requirements and source_coverage_units explicitly.",
            )
        state_manager.add_debug_log(
            f"[PPTXAgent] Content coverage audit: status={audit.status}, "
            f"missing={len(audit.missing_items)}, weak={len(audit.weakly_covered_items)}, "
            f"overcompressed={len(audit.overcompressed_items)}."
        )
        if low_content_slides:
            audit.status = "fail"
            audit.overcompressed_items.append(
                f"Slides with too little visible body content: {low_content_slides}"
            )
            audit.recommendations.append(
                "Fill low-content slides with concrete visible material or remove intentional divider/back-cover slides from the generated body slide list."
            )
        if len(presentation_data.slides) < min_body_slides:
            audit.status = "fail"
            audit.overcompressed_items.append(
                f"Generated {len(presentation_data.slides)} body slides, below required minimum {min_body_slides}."
            )
            audit.recommendations.append(
                "Increase the number of body slides so recommended_storyline sections and source coverage units are not merged too aggressively."
            )
        if audit.status != "pass" and not audit.revision_instruction:
            audit.revision_instruction = (
                "Revise the deck so every missing or weakly covered item is explicitly reflected. "
                f"Use at least {min_body_slides} body slides, add concrete visible text, and populate coverage_refs for every slide."
            )
        return audit

    def _repair_presentation_structure(
        self,
        brief: PresentationSourceBrief,
        current_data: PresentationDSLSchema,
        audits: List[SlideVisualAudit],
        system_instruction: str,
        gen_config: types.GenerateContentConfig,
        model_id: str,
    ) -> PresentationDSLSchema:
        audit_text = _audit_summary_to_text(audits)
        revision_instruction = (
            "VLM監査で以下の問題が検出されました。"
            "レイアウトを読みやすくし、文字量を減らし、不要・不自然な画像生成を削除し、"
            "source brief の重要事実と根拠は維持したまま PresentationDSLSchema を再生成してください。\n"
            "重大な問題がある画像スライドは、メイン_1 などテキスト中心レイアウトへ変更してください。\n"
            "既存構成:\n"
            f"{json.dumps(current_data.model_dump(), ensure_ascii=False, indent=2)[:16000]}\n\n"
            "監査結果:\n"
            f"{audit_text}"
        )
        repaired = self._generate_presentation_structure(
            brief=brief,
            system_instruction=system_instruction,
            gen_config=gen_config,
            model_id=model_id,
            revision_instruction=revision_instruction,
        )
        state_manager.add_debug_log("[PPTXAgent] Presentation structure regenerated from visual audit feedback.")
        return repaired

    def _validate_and_adjust_slides(
        self,
        slides: List[SlideNode],
        layouts_info: Dict[str, Any],
        temp_dir: str,
    ) -> tuple[List[SlideNode], List[int]]:
        offsets = [0] * len(slides)
        for retry_loop in range(3):
            state_manager.add_debug_log(f"[PPTXAgent] Validation Round {retry_loop + 1}")
            for i, slide in enumerate(slides):
                html_content = generate_slide_html(slide, layouts_info, font_size_offset=offsets[i])
                html_path = os.path.join(temp_dir, f"{slide.slide_number:02d}.html")
                with open(html_path, "w", encoding="utf-8") as f:
                    f.write(html_content)
            validation_results = validate_slide_overflow(temp_dir, slides, offsets)
            overflowed_slides = [res for res in validation_results if res.get("overflowed")]
            if not overflowed_slides:
                state_manager.add_debug_log("[PPTXAgent] Geometry validation passed. No overflows detected!")
                break
            if retry_loop == 2:
                state_manager.add_debug_log("[PPTXAgent] Retries exhausted. Applying emergency font size shrink (-2pt) for overflowed slides.")
                for res in overflowed_slides:
                    idx = next(idx for idx, s in enumerate(slides) if s.slide_number == res["slide_number"])
                    offsets[idx] = -2
                break
            state_manager.add_debug_log(f"[PPTXAgent] Overflows detected on {len(overflowed_slides)} slides. Starting Self-Correction...")
            for res in overflowed_slides:
                slide_num = res["slide_number"]
                idx = next(idx for idx, s in enumerate(slides) if s.slide_number == slide_num)
                target_slide = slides[idx]
                scale = res["available_height"] / max(res["required_height"], 1)
                target_ph_idx = int(res["element_id"].split("_")[1])
                target_ph = next((p for p in target_slide.placeholders if p.idx == target_ph_idx), None)
                target_text = target_ph.text_content if target_ph else ""
                retry_prompt = (
                    f"【レイアウトエラー検知】\n"
                    f"スライド番号: {slide_num}\n"
                    f"エラー要素 (プレースホルダー idx): {target_ph_idx}\n"
                    f"現在のテキスト: \"{target_text}\"\n"
                    f"物理的許容高さ: {res['available_height']}px\n"
                    f"要求高さ: {res['required_height']}px\n\n"
                    "上記要素で文字溢れ（Overflow）が発生しました。"
                    f"現在のスライド「{target_slide.title}」のプレースホルダー(idx={target_ph_idx})の箇条書きテキストを、"
                    f"約 {scale:.2f} 倍（文字数として30%〜50%削減）に要約・圧縮してください。"
                )
                correction_config = types.GenerateContentConfig(
                    system_instruction="あなたは短縮・要約のスペシャリストです。指定されたPydanticスキーマに厳密に準拠したJSONを返してください。",
                    max_output_tokens=4096,
                    temperature=0.2,
                    response_mime_type="application/json",
                    response_schema=PlaceholderContent,
                )
                correction_result = llm_router.generate_content_with_route(
                    llm_clients=self.llm_clients,
                    model_id="gemini-3.5-flash",
                    contents=[types.Content(role="user", parts=[types.Part.from_text(text=retry_prompt)])],
                    config=correction_config,
                    mode="pptx_correction",
                    logger=state_manager.add_debug_log,
                )
                new_content = PlaceholderContent.model_validate_json(correction_result.text)
                for p_idx, p_item in enumerate(slides[idx].placeholders):
                    if p_item.idx == target_ph_idx:
                        slides[idx].placeholders[p_idx] = new_content
                        break
        return slides, offsets

    def _prepare_slide_images(
        self,
        slides: List[SlideNode],
        temp_dir: str,
        user_images: Optional[List[dict]],
    ) -> Dict[int, Dict[int, str]]:
        slide_images = {}
        generated_count = 0
        for slide in slides:
            slide_num = slide.slide_number
            slide_images[slide_num] = {}
            for content in slide.placeholders:
                if content.image_prompt:
                    out_path = os.path.join(temp_dir, f"ai_gen_{slide_num}_{content.idx}.png")
                    success = generate_ai_image_with_nano_banana(
                        client=self.llm_clients.standard_client,
                        prompt=content.image_prompt,
                        output_path=out_path,
                    )
                    if success:
                        slide_images[slide_num][content.idx] = out_path
                        generated_count += 1
                elif content.use_user_image and user_images:
                    user_img = user_images[0]
                    img_bytes = user_img["bytes"]
                    img_mime = user_img["type"]
                    out_path = os.path.join(temp_dir, f"user_processed_{slide_num}_{content.idx}.png")
                    if content.crop_instruction:
                        success = crop_user_image_with_llm(
                            client=self.llm_clients.standard_client,
                            image_bytes=img_bytes,
                            image_mime_type=img_mime,
                            instruction=content.crop_instruction,
                            output_path=out_path,
                        )
                        if success:
                            slide_images[slide_num][content.idx] = out_path
                    else:
                        try:
                            with open(out_path, "wb") as f:
                                f.write(img_bytes)
                            slide_images[slide_num][content.idx] = out_path
                        except Exception as e:
                            state_manager.add_debug_log(f"[PPTXAgent] Saving user image failed: {e}", "warning")
        state_manager.add_debug_log(f"[PPTXAgent] Prepared slide images: generated={generated_count}.")
        return slide_images

    def _save_physical_presentation(
        self,
        presentation_data: PresentationDSLSchema,
        slides: List[SlideNode],
        offsets: List[int],
        slide_images: Dict[int, Dict[int, str]],
        final_pptx_path: str,
        template_path: Optional[str],
        has_template: bool,
        layouts_info: Dict[str, Any],
        check_prs: Optional[Presentation],
        source_brief: Optional[PresentationSourceBrief] = None,
    ) -> None:
        if has_template:
            try:
                prs = Presentation(template_path)
            except Exception as e:
                state_manager.add_debug_log(f"[PPTXAgent] Failed to load template, using default blank: {e}", "warning")
                prs = Presentation()
                has_template = False
        else:
            prs = Presentation()
        if not has_template:
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

        def find_layout_by_name(prs_obj, name_str, fallback_idx):
            for layout in prs_obj.slide_layouts:
                if layout.name == name_str:
                    return layout
            if len(prs_obj.slide_layouts) > fallback_idx:
                return prs_obj.slide_layouts[fallback_idx]
            return prs_obj.slide_layouts[0]
        title_layout = find_layout_by_name(prs, "タイトルのみ", 0)
        blank_layout = find_layout_by_name(prs, "コンテンツ 2", 6)
        COLOR_PRIMARY = RGBColor(0x1e, 0x3a, 0x8a)
        COLOR_WHITE = RGBColor(0xff, 0xff, 0xff)
        COLOR_TEXT_MUTED = RGBColor(0x64, 0x74, 0x8b)
        if has_template:
            title_slide = prs.slides[0]
            title_written = False
            for shape in title_slide.shapes:
                if shape.is_placeholder:
                    ph_type = shape.placeholder_format.type
                    if ph_type == 1 or ph_type == 3:
                        shape.text = presentation_data.presentation_title
                        _format_cover_title(shape)
                        title_written = True
                    elif ph_type == 2:
                        shape.text = "R&D Technical Presentation  |  Generated by GP-Chat"
                        for paragraph in shape.text_frame.paragraphs:
                            paragraph.font.name = "Meiryo"
            if not title_written:
                left, top, width, height = _title_box_from_layout(title_slide, prs)
                tx_main_title = title_slide.shapes.add_textbox(left, top, width, height)
                tf_main_title = tx_main_title.text_frame
                tf_main_title.word_wrap = True
                p_main = tf_main_title.paragraphs[0]
                p_main.alignment = PP_ALIGN.CENTER
                p_main.text = presentation_data.presentation_title
                _format_cover_title(tx_main_title)
        else:
            title_slide = prs.slides.add_slide(title_layout)
            bg_title = title_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
            bg_title.fill.solid()
            bg_title.fill.fore_color.rgb = COLOR_WHITE
            bg_title.line.fill.background()
            accent_line = title_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.666), Inches(3.6), Inches(6.0), Inches(0.02))
            accent_line.fill.solid()
            accent_line.fill.fore_color.rgb = COLOR_PRIMARY
            accent_line.line.fill.background()
            tx_main_title = title_slide.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(11.333), Inches(1.1))
            tf_main_title = tx_main_title.text_frame
            tf_main_title.word_wrap = True
            p_main = tf_main_title.paragraphs[0]
            p_main.alignment = PP_ALIGN.CENTER
            p_main.text = presentation_data.presentation_title
            p_main.font.name = "Meiryo"
            p_main.font.size = Pt(44)
            p_main.font.bold = True
            p_main.font.color.rgb = COLOR_PRIMARY
            tx_sub_title = title_slide.shapes.add_textbox(Inches(1.0), Inches(3.9), Inches(11.333), Inches(1.2))
            tf_sub_title = tx_sub_title.text_frame
            tf_sub_title.word_wrap = True
            p_sub = tf_sub_title.paragraphs[0]
            p_sub.alignment = PP_ALIGN.CENTER
            p_sub.text = "R&D Technical Presentation  |  Generated by GP-Chat"
            p_sub.font.name = "Meiryo"
            p_sub.font.size = Pt(18)
            p_sub.font.color.rgb = COLOR_TEXT_MUTED
        for i, slide_data in enumerate(slides):
            layout_to_use = blank_layout
            if has_template and slide_data.layout_name in layouts_info:
                layout_to_use = prs.slide_layouts[layouts_info[slide_data.layout_name]["layout_index"]]
            img_paths_for_slide = slide_images.get(slide_data.slide_number, {})
            render_pptx_slide(
                prs=prs,
                slide_data=slide_data,
                font_size_offset=offsets[i],
                current_index=i,
                total_slides=len(slides),
                image_paths=img_paths_for_slide,
                has_template=has_template,
                blank_layout=layout_to_use,
            )
        if source_brief and source_brief.references:
            if add_references_slide(prs, blank_layout, source_brief.references):
                state_manager.add_debug_log(
                    f"[PPTXAgent] References slide inserted with {min(len(source_brief.references), 8)} visible reference(s)."
                )
        if has_template:
            slide_id_list = prs.slides._sldIdLst

            def remove_slide_at(index: int):
                slide_id = slide_id_list[index]
                slide_id_list.remove(slide_id)
                return slide_id

            def remove_footer_placeholders(slide):
                for shape in list(slide.shapes):
                    if not shape.is_placeholder:
                        continue
                    if shape.placeholder_format.type in (13, 16):
                        sp = shape._element
                        sp.getparent().remove(sp)
            template_slide_count = len(check_prs.slides) if check_prs is not None else 0
            content_layout_names = set(layouts_info.keys())
            back_cover_names = {"end", "裏表紙", "backcover", "back cover", "back_cover"}
            removed_mockups = 0
            back_cover_id = None
            for idx in range(min(template_slide_count, len(prs.slides)) - 1, 0, -1):
                layout_name = prs.slides[idx].slide_layout.name
                normalized_name = layout_name.strip().lower()
                if normalized_name in back_cover_names:
                    remove_footer_placeholders(prs.slides[idx])
                    back_cover_id = remove_slide_at(idx)
                elif layout_name in content_layout_names:
                    remove_slide_at(idx)
                    removed_mockups += 1
            if back_cover_id is not None:
                slide_id_list.append(back_cover_id)
            state_manager.add_debug_log(
                f"[PPTXAgent] Template cleanup: removed {removed_mockups} content mockup slide(s); "
                f"back cover {'moved to end' if back_cover_id is not None else 'not found'}."
            )
        prs.save(final_pptx_path)
        reset_pptx_open_view(final_pptx_path)
        state_manager.add_debug_log(f"[PPTXAgent] PowerPoint file successfully saved: {final_pptx_path}")

    def generate_presentation_pipeline(
        self,
        chat_history: List[dict],
        session_id: str,
        model_id: str = "gemini-3.5-flash",
        user_images: List[dict] = None,
        materialized_contents: Optional[List[Any]] = None,
        materialized_system_instruction: str = "",
        file_attachments_meta: Optional[List[dict]] = None,
        tools_config: Optional[List[Any]] = None,
        conversation_grounding_metadata: Optional[dict] = None,
    ) -> str:
        """会話履歴およびユーザー添付画像リストからPowerPoint資料を作成する同期メインパイプライン。"""
        import streamlit as st
        current_report_folder = st.session_state.get("current_report_folder")
        if not current_report_folder:
            current_chat_filename = st.session_state.get("current_chat_filename")
            if current_chat_filename:
                folder_name = os.path.splitext(os.path.basename(current_chat_filename))[0]
            else:
                folder_name = f"pptx_{session_id[:8]}"
            folder_name = folder_name[:80] or "pptx_chat"
        else:
            folder_name = current_report_folder
        output_dir = os.path.join("slide_data", folder_name)
        os.makedirs(output_dir, exist_ok=True)
        temp_dir = os.path.join("temp_workspace", session_id)
        os.makedirs(temp_dir, exist_ok=True)
        pptx_number = 1
        for filename in os.listdir(output_dir):
            stem, ext = os.path.splitext(filename)
            if ext.lower() == ".pptx" and stem.isdigit():
                pptx_number = max(pptx_number, int(stem) + 1)
        final_pptx_path = os.path.abspath(os.path.join(output_dir, f"{pptx_number:02d}.pptx"))
        # --- テンプレートの自動ロードと解析スキャン ---
        current_dir = os.path.dirname(os.path.abspath(__file__))
        template_pptx = os.path.join(current_dir, "format.pptx")
        template_potx = os.path.join(current_dir, "format.potx")
        template_path = None
        if os.path.exists(template_pptx):
            template_path = template_pptx
        elif os.path.exists(template_potx):
            template_path = template_potx
        layouts_info = {}
        has_template = False
        if template_path:
            try:
                # 3枚以上のスライドがあることを確認
                check_prs = Presentation(template_path)
                if len(check_prs.slides) >= 3:
                    layouts_info = scan_template_layouts(template_path)
                    has_template = True
                    state_manager.add_debug_log(f"[PPTXAgent] Scanned {len(layouts_info)} layouts from format.pptx template.")
                else:
                    state_manager.add_debug_log("[PPTXAgent] Template does not have at least 3 slides.", "warning")
            except Exception as e:
                state_manager.add_debug_log(f"[PPTXAgent] Error inspecting template at start: {e}", "warning")
        # --- 第1層: 構造化 JSON 生成 ---
        state_manager.add_debug_log("[PPTXAgent] Step 1: Requesting Presentation Structure (JSON Schema)")
        # テンプレートレイアウト情報をプロンプトに注入
        template_instruction = ""
        if has_template:
            allowed_layout_names = ", ".join(f"'{name}'" for name in layouts_info.keys())
            template_instruction = (
                "\n\n【利用可能なスライドテンプレートレイアウト情報】\n"
                "以下は本文スライドとして使用できるレイアウトだけを抽出した一覧です。\n"
                f"各スライドの 'layout_name' は必ず次の候補だけから選んでください: {allowed_layout_names}\n"
                "表紙用・裏表紙用・日付・スライド番号だけのレイアウトは本文スライドに使用しないでください。\n"
                "また、'placeholders' の中では、指定したレイアウトの 'placeholders' に定義されている 'idx' だけを指定してください。\n"
                "TITLE/CENTER_TITLE にはスライドタイトル、BODY/CONTENT には本文、IMAGE には image_prompt を入れてください。\n"
                "画像用レイアウトは、構造・関係・プロセス・比較・リスク構図などの解説を視覚的に補強するために積極的に活用してください。\n"
                "抽象的な話題に対しても、概念図（diagram）や比喩表現（metaphor）を用いた説明的なコンセプトイラストを生成させることができます。\n"
                "IMAGE の image_prompt は、文字・ロゴ・ブランド名を含まない、スライド内容を補足・説明する具体的なビジネス向けの図解、コンセプトイラスト、またはシンボル/アイコンを英語で指定してください。\n"
            )
            for name, details in layouts_info.items():
                ph_desc = []
                for ph in details["placeholders"]:
                    ph_desc.append(f"  * idx={ph['idx']} (タイプ: {ph['type']}, 名前: '{ph['name']}')")
                template_instruction += f"- レイアウト名: '{name}'\n  使用可能な入力枠 (placeholders):\n" + "\n".join(ph_desc) + "\n"
        else:
            template_instruction = (
                "\n\n(テンプレートがロードされなかったため、デフォルトの白紙レイアウトで生成します。)\n"
                "layout_name には 'blank' などの仮の文字列を入れ、placeholders は idx=0 などのダミーのテキストを流し込んでください。"
            )
        source_system_instruction = ""
        if materialized_system_instruction:
            source_system_instruction = (
                "\n\n【元チャットのシステム/モード指示】\n"
                "以下は元の会話で使われた前提指示です。資料化の目的・制約・分析方針として尊重してください。\n"
                f"{materialized_system_instruction[:6000]}"
            )
        system_instruction = (
            "あなたはPowerPoint資料の構造化プロのデザイナーです。\n"
            "これまでの会話履歴を参照して、最終報告書やプレゼンテーション用の構成ストーリーを作成してください。\n"
            "添付ファイル・添付画像・検索結果・キャンバスが会話コンテキストに含まれている場合は、必ず主要な根拠として読み込み、"
            "テンプレート文言だけで資料を作ってはいけません。\n"
            "提示された PresentationDSLSchema に従って、論理的に厳格で情報の引き算がなされたスライド構成をJSONとして出力してください。\n"
            "スライド枚数は3〜5枚が適切です。\n"
            "単なる箇条書きの羅列ではなく、1枚ごとに「結論→根拠→示唆」が読めるインフォグラフィック風の構成にしてください。\n"
            "比較、時系列、因果関係、リスク/対策、意思決定ポイントなど、視覚的に整理しやすい切り口を優先してください。\n"
            "BODY には1〜2行のキーメッセージだけを置き、CONTENT には短い見出し付きの要点を2〜3個に絞ってください。\n"
            "各箇条書きは文字溢れを防ぐため原則30文字以内で要約してください。\n"
            "スライドの理解度を深めるため、説明価値が高い場合は積極的に画像生成（IMAGE プレースホルダー）を取り入れてください（目安として全体で最大4枚程度まで）。\n"
            "スライド全体のタイトルも placeholders の中に適切な TITLE または CENTER_TITLE タイプがある場合は必ず含めて出力してください。"
            + source_system_instruction
            + template_instruction
        )
        system_instruction += (
            "\n\n[Content coverage policy]\n"
            "The deck is a report-generation output, not a teaser. Use all available relevant information from the source brief.\n"
            "Do not discard concrete facts, dates, numbers, filenames, user constraints, image observations, or web evidence merely to keep the deck short.\n"
            "If the source brief contains many important facts, add slides and organize them into sections instead of overcompressing.\n"
            "Every slide must set coverage_refs to trace which brief items it covers.\n"
            "Use uncertainty notes when evidence is weak; do not turn weak evidence into a strong claim.\n"
            "\n\n[Generic visual_type guidance]\n"
            "For each slide, set visual_type by information structure, not by subject/domain.\n"
            "- timeline: dates, phases, milestones, sequence over time.\n"
            "- process: steps, workflow, input/output, approval or operational flow.\n"
            "- comparison: multiple options/entities with evaluation axes.\n"
            "- kpi: numeric facts, metrics, quantities, rates, costs, counts.\n"
            "- matrix: two-axis prioritization, importance/urgency, impact/probability.\n"
            "- risk: risk/issue/cause/impact/mitigation structure.\n"
            "- summary: grouped takeaways without a stronger structure.\n"
            "- none: use only when an added visual would reduce clarity.\n"
            "Also set visual_variant to avoid repetitive slides. Valid useful variants include:\n"
            "- timeline: vertical_timeline, phase_bands, horizontal_timeline, milestone_cards, now_next_later, gantt_roadmap.\n"
            "- process: vertical_flow, chevron_flow, numbered_steps, loop_cycle, swimlane_flow, funnel.\n"
            "- comparison: cards_2x2, scorecard, pros_cons, table_compare, ranked_bars, before_after, option_columns.\n"
            "- kpi: metric_cards, big_numbers, scorecard, progress_bars, gauge_cards, delta_callouts, waterfall.\n"
            "- matrix: risk_matrix, priority_quadrants, decision_matrix.\n"
            "- risk: risk_matrix, cause_impact_mitigation, risk_register, heatmap, escalation_ladder.\n"
            "- summary: bands, pyramid, hub_spoke.\n"
            "Use auto when unsure, but vary variants across adjacent slides when the information structure allows it.\n"
            "Keep the template layout names and placeholder idx values valid. Do not rely on topic-specific words.\n"
        )
        gen_config = types.GenerateContentConfig(
            system_instruction=system_instruction,
            max_output_tokens=65536,
            temperature=0.2,
            response_mime_type="application/json",
            response_schema=PresentationDSLSchema,
        )
        if materialized_contents:
            base_chat_contents = [_clone_content_for_pptx(content) for content in materialized_contents]
            state_manager.add_debug_log(
                f"[PPTXAgent] Using materialized multimodal context: {len(base_chat_contents)} content item(s)."
            )
        else:
            base_chat_contents = []
            for msg in chat_history:
                if msg["role"] == "system":
                    continue
                base_chat_contents.append(
                    types.Content(
                        role=_normalize_api_role(msg.get("role", "user")),
                        parts=[types.Part.from_text(text=msg.get("content", ""))]
                    )
                )
            state_manager.add_debug_log(
                f"[PPTXAgent] Using plain chat history only: {len(base_chat_contents)} content item(s).",
                "warning",
            )
        attachment_summary = _format_attachment_summary(file_attachments_meta)
        research_context, research_grounding_metadata = self._run_research_context(
            chat_contents=base_chat_contents,
            tools_config=tools_config,
            model_id=model_id,
        )
        combined_grounding_metadata = research_grounding_metadata
        existing_grounding_text = _format_grounding_metadata(conversation_grounding_metadata)
        if existing_grounding_text:
            research_context = (
                research_context
                + "\n\n【会話内の既存検索メタデータ】\n"
                + existing_grounding_text
            )
            combined_grounding_metadata = llm_router.merge_grounding_metadata(
                combined_grounding_metadata,
                conversation_grounding_metadata,
            )
            state_manager.add_debug_log("[PPTXAgent] Added existing conversation grounding metadata to source brief.")
        conversation_excerpt = _build_conversation_excerpt(chat_history)
        source_brief = self._build_source_brief(
            chat_contents=base_chat_contents,
            conversation_excerpt=conversation_excerpt,
            attachment_summary=attachment_summary,
            file_attachments_meta=file_attachments_meta,
            research_context=research_context,
            grounding_metadata=combined_grounding_metadata,
            materialized_system_instruction=materialized_system_instruction,
            model_id=model_id,
        )
        try:
            brief_path = os.path.join(output_dir, f"{pptx_number:02d}_source_brief.json")
            with open(brief_path, "w", encoding="utf-8") as f:
                json.dump(source_brief.model_dump(), f, ensure_ascii=False, indent=2)
            state_manager.add_debug_log(f"[PPTXAgent] Source brief saved: {brief_path}")
        except Exception as e:
            state_manager.add_debug_log(f"[PPTXAgent] Source brief save skipped: {e}", "warning")
        presentation_data = self._generate_presentation_structure(
            brief=source_brief,
            system_instruction=system_instruction,
            gen_config=gen_config,
            model_id=model_id,
        )
        _attach_reference_usage(source_brief, presentation_data)
        coverage_audit = self._audit_content_coverage(
            brief=source_brief,
            presentation_data=presentation_data,
            model_id=model_id,
        )
        if coverage_audit.status != "pass":
            state_manager.add_debug_log(
                "[PPTXAgent] Content coverage audit found gaps. Regenerating structure once.",
                "warning",
            )
            missing_summary = {
                "missing_items": coverage_audit.missing_items,
                "weakly_covered_items": coverage_audit.weakly_covered_items,
                "overcompressed_items": coverage_audit.overcompressed_items,
                "recommendations": coverage_audit.recommendations,
            }
            revision_instruction = (
                "Content coverage audit found gaps before PPTX rendering.\n"
                "Fix the structure so all missing and weakly covered items are explicitly represented.\n"
                "Add slides when needed. Do not compress distinct facts into generic statements.\n"
                "Every slide must include coverage_refs that quote the source brief items it covers, but coverage_refs alone are not enough.\n"
                "The covered information must also appear in visible placeholder text or image prompts.\n"
                "Do not create empty divider/back-cover slides in the generated body slide list.\n\n"
                f"Audit details:\n{json.dumps(missing_summary, ensure_ascii=False, indent=2)}\n\n"
                f"Audit instruction:\n{coverage_audit.revision_instruction}"
            )
            presentation_data = self._generate_presentation_structure(
                brief=source_brief,
                system_instruction=system_instruction,
                gen_config=gen_config,
                model_id=model_id,
                revision_instruction=revision_instruction,
            )
            _attach_reference_usage(source_brief, presentation_data)
            coverage_audit = self._audit_content_coverage(
                brief=source_brief,
                presentation_data=presentation_data,
                model_id=model_id,
            )
        try:
            brief_path = os.path.join(output_dir, f"{pptx_number:02d}_source_brief.json")
            with open(brief_path, "w", encoding="utf-8") as f:
                json.dump(source_brief.model_dump(), f, ensure_ascii=False, indent=2)
            state_manager.add_debug_log(f"[PPTXAgent] Source brief updated with reference usage: {brief_path}")
        except Exception as e:
            state_manager.add_debug_log(f"[PPTXAgent] Source brief reference update skipped: {e}", "warning")
        try:
            structure_path = os.path.join(output_dir, f"{pptx_number:02d}_presentation_structure.json")
            with open(structure_path, "w", encoding="utf-8") as f:
                json.dump(presentation_data.model_dump(), f, ensure_ascii=False, indent=2)
            state_manager.add_debug_log(f"[PPTXAgent] Presentation structure saved: {structure_path}")
        except Exception as e:
            state_manager.add_debug_log(f"[PPTXAgent] Presentation structure save skipped: {e}", "warning")
        try:
            coverage_path = os.path.join(output_dir, f"{pptx_number:02d}_coverage_audit.json")
            with open(coverage_path, "w", encoding="utf-8") as f:
                json.dump(coverage_audit.model_dump(), f, ensure_ascii=False, indent=2)
            state_manager.add_debug_log(f"[PPTXAgent] Content coverage audit saved: {coverage_path}")
        except Exception as e:
            state_manager.add_debug_log(f"[PPTXAgent] Content coverage audit save skipped: {e}", "warning")
        try:
            references_path = os.path.join(output_dir, f"{pptx_number:02d}_references.json")
            with open(references_path, "w", encoding="utf-8") as f:
                json.dump([item.model_dump() for item in source_brief.references], f, ensure_ascii=False, indent=2)
            state_manager.add_debug_log(f"[PPTXAgent] References saved: {references_path}")
        except Exception as e:
            state_manager.add_debug_log(f"[PPTXAgent] References save skipped: {e}", "warning")
        # --- 第2層以降: 幾何学検証、画像準備、PPTX生成、VLM監査 ---
        state_manager.add_debug_log("[PPTXAgent] Step 2-4: Validating layout and generating physical PowerPoint presentation")
        slides, offsets = self._validate_and_adjust_slides(
            slides=presentation_data.slides,
            layouts_info=layouts_info,
            temp_dir=temp_dir,
        )
        slide_images = self._prepare_slide_images(
            slides=slides,
            temp_dir=temp_dir,
            user_images=user_images,
        )
        self._save_physical_presentation(
            presentation_data=presentation_data,
            slides=slides,
            offsets=offsets,
            slide_images=slide_images,
            final_pptx_path=final_pptx_path,
            template_path=template_path,
            has_template=has_template,
            layouts_info=layouts_info,
            check_prs=check_prs if "check_prs" in locals() else None,
            source_brief=source_brief,
        )
        try:
            audits = run_visual_quality_audit(
                client=self.llm_clients.standard_client,
                pptx_path=final_pptx_path,
                temp_dir=temp_dir,
                model_id=model_id,
            )
            if any(audit.status == "fail" for audit in audits):
                state_manager.add_debug_log("[PPTXAgent] Visual audit found serious layout issues. Regenerating once from audit feedback.", "warning")
                presentation_data = self._repair_presentation_structure(
                    brief=source_brief,
                    current_data=presentation_data,
                    audits=audits,
                    system_instruction=system_instruction,
                    gen_config=gen_config,
                    model_id=model_id,
                )
                _attach_reference_usage(source_brief, presentation_data)
                slides, offsets = self._validate_and_adjust_slides(
                    slides=presentation_data.slides,
                    layouts_info=layouts_info,
                    temp_dir=temp_dir,
                )
                slide_images = self._prepare_slide_images(
                    slides=slides,
                    temp_dir=temp_dir,
                    user_images=user_images,
                )
                self._save_physical_presentation(
                    presentation_data=presentation_data,
                    slides=slides,
                    offsets=offsets,
                    slide_images=slide_images,
                    final_pptx_path=final_pptx_path,
                    template_path=template_path,
                    has_template=has_template,
                    layouts_info=layouts_info,
                    check_prs=check_prs if "check_prs" in locals() else None,
                    source_brief=source_brief,
                )
                final_audits = run_visual_quality_audit(
                    client=self.llm_clients.standard_client,
                    pptx_path=final_pptx_path,
                    temp_dir=temp_dir,
                    model_id=model_id,
                )
                if final_audits:
                    failed = sum(1 for audit in final_audits if audit.status == "fail")
                    warned = sum(1 for audit in final_audits if audit.status == "warning")
                    state_manager.add_debug_log(
                        f"[PPTXAgent] Final visual audit after repair: {failed} fail, {warned} warning."
                    )
        except Exception as e:
            state_manager.add_debug_log(f"[PPTXAgent] Visual audit/repair skipped due to error: {e}", "warning")
        try:
            shutil.rmtree(temp_dir)
        except Exception as e:
            state_manager.add_debug_log(f"[PPTXAgent] Non-critical cleanup warning: {e}", "warning")
        return final_pptx_path
